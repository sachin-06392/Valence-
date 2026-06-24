from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, StreamingResponse
from pydantic import BaseModel
from typing import Optional, Any, Dict, List
from functools import lru_cache
import os
import re
import tempfile
import io
import requests
import yfinance as yf
import numpy as np
from datetime import datetime, timezone

from company_universe import COMPANY_UNIVERSE
from financials_db import FINANCIALS_DB
from report_generator import generate_banker_report


app = FastAPI(title="CompFinder API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "http://localhost:3000",
        "http://127.0.0.1:3000",
        "http://localhost:5173",
        "http://127.0.0.1:5173",
        "https://valence-lac-ten.vercel.app",
    ],
    allow_origin_regex=r"https://.*\.vercel\.app",
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# -----------------------------------------------------------------------------
# Sector taxonomy / public comparable universe
# -----------------------------------------------------------------------------
# The app used to have 5 broad buckets. This version uses a more banker-style
# software / tech taxonomy and separates direct operating comps from strategic
# ecosystem comps.

LIVE = "live"
HISTORICAL = "historical"
DIRECT = "direct"
STRATEGIC = "strategic"


def comp(ticker: str, name: str, status: str = LIVE, note: str = "", yf_symbol: Optional[str] = None):
    return {
        "ticker": ticker.upper(),
        "symbol": ticker.upper(),
        "name": name,
        "companyName": name,
        "status": status,
        "note": note,
        "yf_symbol": yf_symbol or ticker.upper(),
    }


SECTOR_TAXONOMY = {
    "AI / Machine Learning": [
        comp("PLTR", "Palantir"),
        comp("SNOW", "Snowflake"),
        comp("DDOG", "Datadog"),
        comp("MDB", "MongoDB"),
        comp("CFLT", "Confluent"),
        comp("ESTC", "Elastic"),
        comp("AI", "C3.ai"),
        comp("PATH", "UiPath"),
    ],
    "Enterprise Software / SaaS": [
        comp("CRM", "Salesforce"),
        comp("NOW", "ServiceNow"),
        comp("WDAY", "Workday"),
        comp("TEAM", "Atlassian"),
        comp("ADBE", "Adobe"),
        comp("INTU", "Intuit"),
        comp("ORCL", "Oracle"),
        comp("SAP", "SAP"),
        comp("MSFT", "Microsoft"),
    ],
    "Cloud / Data Infrastructure": [
        comp("SNOW", "Snowflake"),
        comp("MDB", "MongoDB"),
        comp("NET", "Cloudflare"),
        comp("DDOG", "Datadog"),
        comp("CFLT", "Confluent"),
        comp("ESTC", "Elastic"),
        comp("FIVN", "Five9"),
        comp("TWLO", "Twilio"),
        comp("DOCN", "DigitalOcean"),
    ],
    "Cybersecurity": [
        comp("CRWD", "CrowdStrike"),
        comp("ZS", "Zscaler"),
        comp("PANW", "Palo Alto Networks"),
        comp("FTNT", "Fortinet"),
        comp("OKTA", "Okta"),
        comp("S", "SentinelOne"),
        comp("CYBR", "CyberArk"),
        comp("TENB", "Tenable"),
        comp("QLYS", "Qualys"),
        comp("NET", "Cloudflare"),
    ],
    "Developer Tools / DevOps": [
        comp("GTLB", "GitLab"),
        comp("TEAM", "Atlassian"),
        comp("DDOG", "Datadog"),
        comp("FROG", "JFrog"),
        comp("CFLT", "Confluent"),
        comp("MDB", "MongoDB"),
        comp("ESTC", "Elastic"),
        comp("DOCN", "DigitalOcean"),
    ],
    "Data Analytics / Observability": [
        comp("DDOG", "Datadog"),
        comp("PLTR", "Palantir"),
        comp("SNOW", "Snowflake"),
        comp("ESTC", "Elastic"),
        comp("MDB", "MongoDB"),
        comp("CFLT", "Confluent"),
        comp("SPLK", "Splunk", HISTORICAL, "Historical comp only; acquired by Cisco in 2024."),
        comp("NEWR", "New Relic", HISTORICAL, "Historical comp only; taken private in 2023."),
    ],
    "FinTech / Payments": [
        comp("PYPL", "PayPal"),
        comp("SQ", "Block"),
        comp("FI", "Fiserv"),
        comp("FIS", "Fidelity National Information Services"),
        comp("GPN", "Global Payments"),
        comp("TOST", "Toast"),
        comp("FOUR", "Shift4 Payments"),
        comp("ADYEN", "Adyen", yf_symbol="ADYEN.AS"),
        comp("NU", "Nu Holdings"),
        comp("SOFI", "SoFi"),
    ],
    "E-Commerce / Marketplaces": [
        comp("SHOP", "Shopify"),
        comp("AMZN", "Amazon"),
        comp("ETSY", "Etsy"),
        comp("EBAY", "eBay"),
        comp("MELI", "MercadoLibre"),
        comp("SE", "Sea Limited"),
        comp("CPNG", "Coupang"),
        comp("DASH", "DoorDash"),
        comp("UBER", "Uber"),
        comp("ABNB", "Airbnb"),
    ],
    "Consumer Internet / Digital Ads": [
        comp("META", "Meta"),
        comp("GOOGL", "Alphabet"),
        comp("SNAP", "Snap"),
        comp("PINS", "Pinterest"),
        comp("RDDT", "Reddit"),
        comp("MTCH", "Match Group"),
        comp("SPOT", "Spotify"),
        comp("NFLX", "Netflix"),
    ],
    "Media / Communications": [
        comp("ZM", "Zoom"),
        comp("TWLO", "Twilio"),
        comp("FIVN", "Five9"),
        comp("NFLX", "Netflix"),
        comp("SPOT", "Spotify"),
        comp("DIS", "Disney"),
        comp("WBD", "Warner Bros. Discovery"),
        comp("PARA", "Paramount"),
    ],
    "Semiconductors / AI Hardware": [
        comp("NVDA", "NVIDIA"),
        comp("AMD", "Advanced Micro Devices"),
        comp("AVGO", "Broadcom"),
        comp("INTC", "Intel"),
        comp("QCOM", "Qualcomm"),
        comp("MRVL", "Marvell"),
        comp("ARM", "Arm Holdings"),
        comp("TSM", "Taiwan Semiconductor"),
        comp("MU", "Micron"),
        comp("ASML", "ASML"),
    ],
    "IT Services / Consulting": [
        comp("ACN", "Accenture"),
        comp("IBM", "IBM"),
        comp("CTSH", "Cognizant"),
        comp("INFY", "Infosys"),
        comp("WIT", "Wipro"),
        comp("GLOB", "Globant"),
        comp("EPAM", "EPAM Systems"),
        comp("DXC", "DXC Technology"),
    ],
    "Healthcare Technology": [
        comp("VEEV", "Veeva"),
        comp("TDOC", "Teladoc"),
        comp("DOCS", "Doximity"),
        comp("HIMS", "Hims & Hers"),
        comp("OSCR", "Oscar Health"),
        comp("CLOV", "Clover Health"),
        comp("HQY", "HealthEquity"),
        comp("OMCL", "Omnicell"),
    ],
    "Life Sciences Software": [
        comp("VEEV", "Veeva"),
        comp("IQV", "IQVIA"),
        comp("CERT", "Certara"),
        comp("MEDP", "Medpace"),
        comp("DHR", "Danaher"),
        comp("TMO", "Thermo Fisher"),
        comp("A", "Agilent"),
    ],
    "HR / Payroll / Workforce Software": [
        comp("WDAY", "Workday"),
        comp("ADP", "ADP"),
        comp("PAYX", "Paychex"),
        comp("PCTY", "Paylocity"),
        comp("PAYC", "Paycom"),
        comp("DAY", "Dayforce"),
        comp("BILL", "BILL"),
    ],
    "Vertical Software": [
        comp("VEEV", "Veeva"),
        comp("APPF", "AppFolio"),
        comp("LSPD", "Lightspeed"),
        comp("TOST", "Toast"),
        comp("MNDY", "Monday.com"),
        comp("HUBS", "HubSpot"),
        comp("DOCU", "DocuSign"),
        comp("BL", "BlackLine"),
        comp("NCNO", "nCino"),
    ],
    "Consumer / SMB Software": [
        comp("INTU", "Intuit"),
        comp("HUBS", "HubSpot"),
        comp("BILL", "BILL"),
        comp("DOCU", "DocuSign"),
        comp("BOX", "Box"),
        comp("DBX", "Dropbox"),
        comp("ASAN", "Asana"),
        comp("MNDY", "Monday.com"),
        comp("SMAR", "Smartsheet", HISTORICAL, "Historical comp only; acquisition closed in 2025."),
    ],
    "Real Estate / PropTech": [
        comp("Z", "Zillow"),
        comp("OPEN", "Opendoor"),
        comp("RDFN", "Redfin"),
        comp("APPF", "AppFolio"),
        comp("CSGP", "CoStar"),
        comp("CBRE", "CBRE"),
        comp("JLL", "JLL"),
    ],
    "Education Technology": [
        comp("DUOL", "Duolingo"),
        comp("COUR", "Coursera"),
        comp("CHGG", "Chegg"),
        comp("LRN", "Stride"),
        comp("TWOU", "2U", HISTORICAL, "Distressed / historical comp only; Chapter 11 and delisting issues."),
        comp("UDMY", "Udemy"),
    ],
    "Energy / Climate Tech": [
        comp("TSLA", "Tesla"),
        comp("ENPH", "Enphase"),
        comp("SEDG", "SolarEdge"),
        comp("FSLR", "First Solar"),
        comp("BE", "Bloom Energy"),
        comp("PLUG", "Plug Power"),
        comp("RUN", "Sunrun"),
        comp("NEE", "NextEra Energy"),
    ],
}


SECTOR_LIST = list(SECTOR_TAXONOMY.keys())

# Backward compatibility for old frontend options / old company_universe sectors.
LEGACY_SECTOR_MAP = {
    "Enterprise SaaS / Tech": "Enterprise Software / SaaS",
    "Healthcare / Biotech": "Healthcare Technology",
    "Financial Services": "FinTech / Payments",
    "Consumer / Retail": "E-Commerce / Marketplaces",
    "Industrials / Manufacturing": "IT Services / Consulting",
    "saas_tech": "Enterprise Software / SaaS",
    "healthcare": "Healthcare Technology",
    "financial": "FinTech / Payments",
    "consumer": "E-Commerce / Marketplaces",
    "industrials": "IT Services / Consulting",
}

INPUT_SECTOR_MAP = {sector: sector for sector in SECTOR_LIST}
INPUT_SECTOR_MAP.update(LEGACY_SECTOR_MAP)
SECTOR_LABEL_MAP = {v: v for v in SECTOR_LIST}
SECTOR_LABEL_MAP.update({k: v for k, v in LEGACY_SECTOR_MAP.items()})

AI_STRATEGIC_ECOSYSTEM_SECTORS = [
    "Semiconductors / AI Hardware",
    "Cloud / Data Infrastructure",
]


SECTOR_KEYWORDS = {
    "AI / Machine Learning": [
        "ai",
        "artificial intelligence",
        "machine learning",
        "llm",
        "large language model",
        "foundation model",
        "generative ai",
        "anthropic",
        "openai",
        "model",
        "agent",
        "copilot",
    ],
    "Enterprise Software / SaaS": [
        "saas",
        "enterprise software",
        "subscription",
        "workflow",
        "crm",
        "erp",
        "business software",
    ],
    "Cloud / Data Infrastructure": [
        "cloud",
        "database",
        "data infrastructure",
        "data warehouse",
        "api",
        "storage",
        "compute",
        "infrastructure",
    ],
    "Cybersecurity": [
        "security",
        "cybersecurity",
        "identity",
        "endpoint",
        "zero trust",
        "firewall",
        "threat",
        "vulnerability",
    ],
    "Developer Tools / DevOps": [
        "developer",
        "devops",
        "git",
        "ci/cd",
        "deployment",
        "software development",
        "engineering workflow",
    ],
    "Data Analytics / Observability": [
        "analytics",
        "observability",
        "monitoring",
        "logs",
        "metrics",
        "business intelligence",
    ],
    "FinTech / Payments": [
        "fintech",
        "payments",
        "banking",
        "lending",
        "card",
        "wallet",
        "merchant",
        "financial services",
    ],
    "E-Commerce / Marketplaces": [
        "ecommerce",
        "e-commerce",
        "marketplace",
        "commerce",
        "retail",
        "online store",
        "delivery",
        "booking",
    ],
    "Consumer Internet / Digital Ads": [
        "consumer internet",
        "ads",
        "advertising",
        "social",
        "search",
        "creator",
        "streaming",
        "dating",
    ],
    "Media / Communications": [
        "media",
        "communications",
        "video",
        "streaming",
        "messaging",
        "contact center",
        "telephony",
    ],
    "Semiconductors / AI Hardware": [
        "semiconductor",
        "chip",
        "gpu",
        "ai hardware",
        "datacenter hardware",
        "foundry",
        "memory",
    ],
    "IT Services / Consulting": [
        "consulting",
        "it services",
        "systems integration",
        "outsourcing",
        "digital transformation",
    ],
    "Healthcare Technology": [
        "healthcare",
        "health tech",
        "telehealth",
        "patient",
        "provider",
        "insurance",
        "clinical",
    ],
    "Life Sciences Software": [
        "life sciences",
        "pharma",
        "clinical trial",
        "biotech",
        "laboratory",
        "drug development",
    ],
    "HR / Payroll / Workforce Software": [
        "hr",
        "payroll",
        "workforce",
        "human capital",
        "benefits",
        "employee",
    ],
    "Vertical Software": [
        "vertical software",
        "industry software",
        "restaurant software",
        "legal software",
        "real estate software",
        "field service",
    ],
    "Consumer / SMB Software": [
        "smb",
        "small business",
        "consumer software",
        "productivity",
        "collaboration",
        "documents",
        "invoicing",
    ],
    "Real Estate / PropTech": [
        "real estate",
        "property",
        "proptech",
        "brokerage",
        "housing",
        "rentals",
    ],
    "Education Technology": [
        "education",
        "edtech",
        "learning",
        "courses",
        "students",
        "tutoring",
        "online education",
    ],
    "Energy / Climate Tech": [
        "energy",
        "climate",
        "solar",
        "renewable",
        "battery",
        "ev",
        "clean energy",
    ],
}


def clean_text(value: Any) -> str:
    if value is None:
        return ""
    return str(value).lower().strip()


def safe_number(value: Any, default=None):
    if value is None or value == "":
        return default

    if isinstance(value, (int, float)):
        return float(value) if np.isfinite(value) else default

    try:
        text = str(value).strip()
        if not text or text in {"—", "-", "N/A", "n/a", "None"}:
            return default

        multiplier = 1
        if text[-1:].upper() == "B":
            multiplier = 1_000_000_000
            text = text[:-1]
        elif text[-1:].upper() == "M":
            multiplier = 1_000_000
            text = text[:-1]

        cleaned = re.sub(r"[$,%x×,\s]", "", text)
        if not cleaned:
            return default

        return float(cleaned) * multiplier
    except Exception:
        return default


def safe_divide(numerator: Any, denominator: Any, default=None):
    num = safe_number(numerator)
    den = safe_number(denominator)

    if num is None or den in [None, 0]:
        return default

    return num / den


def safe_report_name(value: Any, default: str = "company") -> str:
    cleaned = re.sub(r"[^A-Za-z0-9_-]+", "", str(value or default).upper())
    return cleaned or default


def normalize_sector(sector: Optional[str]) -> Optional[str]:
    if not sector:
        return None

    raw = str(sector).strip()

    if raw in INPUT_SECTOR_MAP:
        return INPUT_SECTOR_MAP[raw]

    raw_lower = raw.lower()

    for valid_sector in SECTOR_LIST:
        if raw_lower == valid_sector.lower():
            return valid_sector

    for old_sector, new_sector in LEGACY_SECTOR_MAP.items():
        if raw_lower == old_sector.lower():
            return new_sector

    for valid_sector in SECTOR_LIST:
        simple = valid_sector.lower().replace("/", " ").replace("-", " ")
        if raw_lower in simple or simple in raw_lower:
            return valid_sector

    return None


def infer_sector(priv: Dict[str, Any]) -> str:
    explicit_sector = normalize_sector(
        priv.get("sector")
        or priv.get("industry")
        or priv.get("category")
        or priv.get("sub_sector")
        or priv.get("subSector")
    )

    if explicit_sector:
        return explicit_sector

    searchable_text = " ".join(
        [
            clean_text(priv.get("company_name")),
            clean_text(priv.get("companyName")),
            clean_text(priv.get("name")),
            clean_text(priv.get("business_description")),
            clean_text(priv.get("businessDescription")),
            clean_text(priv.get("description")),
            clean_text(priv.get("revenue_model")),
            clean_text(priv.get("revenueModel")),
            clean_text(priv.get("customer_type")),
            clean_text(priv.get("customerType")),
        ]
    )

    scores = {}

    for sector, keywords in SECTOR_KEYWORDS.items():
        score = 0
        for keyword in keywords:
            if keyword in searchable_text:
                score += 1
        if score > 0:
            scores[sector] = score

    if not scores:
        return "Enterprise Software / SaaS"

    return max(scores.items(), key=lambda item: item[1])[0]


def candidate_reason(selected_sector: str, candidate_sector: str, relationship: str, candidate: Dict[str, Any]) -> str:
    if candidate.get("status") == HISTORICAL:
        return candidate.get("note") or "Historical comp only."

    if relationship == STRATEGIC:
        return (
            f"Strategic ecosystem comp for {selected_sector}; useful for AI infrastructure context, "
            "but excluded from direct operating valuation multiples."
        )

    return f"Direct operating public comp in {candidate_sector}."


def build_candidate_pool(selected_sector: str, include_historical: bool = False) -> List[Dict[str, Any]]:
    pool = []
    seen = set()

    def add_sector(sector_name: str, relationship: str):
        for order_index, candidate in enumerate(SECTOR_TAXONOMY.get(sector_name, [])):
            if candidate.get("status") == HISTORICAL and not include_historical:
                continue

            ticker = candidate["ticker"].upper()
            dedupe_key = (ticker, relationship)

            if dedupe_key in seen:
                continue

            seen.add(dedupe_key)
            pool.append(
                {
                    **candidate,
                    "sector": sector_name,
                    "industry": sector_name,
                    "sub": sector_name,
                    "subSector": sector_name,
                    "relationship": relationship,
                    "compType": relationship,
                    "order_index": order_index,
                    "reason": candidate_reason(selected_sector, sector_name, relationship, candidate),
                }
            )

    add_sector(selected_sector, DIRECT)

    # For AI companies like Anthropic, show AI infrastructure names separately.
    # These are strategic ecosystem comps, not direct operating comps.
    if selected_sector == "AI / Machine Learning":
        for ecosystem_sector in AI_STRATEGIC_ECOSYSTEM_SECTORS:
            add_sector(ecosystem_sector, STRATEGIC)

    return pool


def get_company_universe_meta(ticker: str) -> Dict[str, Any]:
    ticker = str(ticker or "").upper()

    for company in COMPANY_UNIVERSE:
        if str(company.get("ticker", "")).upper() == ticker:
            return company

    return {}


def normalize_financial_record(fin: Dict[str, Any]) -> Dict[str, Any]:
    if not isinstance(fin, dict):
        return {}

    revenue = safe_number(fin.get("revenue"))
    ebitda = safe_number(fin.get("ebitda"))
    gross_profit = safe_number(fin.get("gross_profit") or fin.get("grossProfit"))
    ev = safe_number(fin.get("ev") or fin.get("enterpriseValue"))
    market_cap = safe_number(fin.get("market_cap") or fin.get("marketCap"))

    ev_rev = safe_number(fin.get("ev_rev") or fin.get("evRevenue"))
    ev_ebitda = safe_number(fin.get("ev_ebitda") or fin.get("evEbitda"))
    ev_gp = safe_number(fin.get("ev_gp") or fin.get("evGrossProfit"))
    pe_ratio = safe_number(fin.get("pe_ratio") or fin.get("peRatio"))

    if ev_rev is None and ev and revenue and revenue > 0:
        ev_rev = ev / revenue

    if ev_ebitda is None and ev and ebitda and ebitda > 0:
        ev_ebitda = ev / ebitda

    if ev_gp is None and ev and gross_profit and gross_profit > 0:
        ev_gp = ev / gross_profit

    rev_growth = safe_number(fin.get("rev_growth") or fin.get("revenueGrowth"))
    gross_margin = safe_number(fin.get("gross_margin") or fin.get("grossMargin"))

    # If values came in as percentages instead of decimals, convert to decimals.
    if rev_growth is not None and abs(rev_growth) > 2:
        rev_growth = rev_growth / 100

    if gross_margin is not None and abs(gross_margin) > 2:
        gross_margin = gross_margin / 100

    return {
        **fin,
        "revenue": revenue,
        "ebitda": ebitda,
        "gross_profit": gross_profit,
        "ev": ev,
        "market_cap": market_cap,
        "ev_rev": ev_rev,
        "ev_ebitda": ev_ebitda,
        "ev_gp": ev_gp,
        "pe_ratio": pe_ratio,
        "rev_growth": rev_growth,
        "gross_margin": gross_margin,
        "employees": fin.get("employees"),
    }


@lru_cache(maxsize=512)
def fetch_yfinance_financials(display_ticker: str, yf_symbol: Optional[str] = None) -> Dict[str, Any]:
    symbol = yf_symbol or display_ticker

    try:
        stock = yf.Ticker(symbol)
        info = stock.info or {}

        revenue = safe_number(info.get("totalRevenue"))
        ebitda = safe_number(info.get("ebitda"))
        gross_profit = safe_number(info.get("grossProfits"))
        ev = safe_number(info.get("enterpriseValue"))
        market_cap = safe_number(info.get("marketCap"))

        ev_rev = safe_number(info.get("enterpriseToRevenue"))
        ev_ebitda = safe_number(info.get("enterpriseToEbitda"))
        pe_ratio = safe_number(info.get("trailingPE") or info.get("forwardPE"))
        rev_growth = safe_number(info.get("revenueGrowth"))
        gross_margin = safe_number(info.get("grossMargins"))

        ev_gp = None
        if ev and gross_profit and gross_profit > 0:
            ev_gp = ev / gross_profit

        if ev_rev is None and ev and revenue and revenue > 0:
            ev_rev = ev / revenue

        if ev_ebitda is None and ev and ebitda and ebitda > 0:
            ev_ebitda = ev / ebitda

        if not any([revenue, ev, market_cap, ev_rev, ev_ebitda]):
            return {}

        return normalize_financial_record(
            {
                "revenue": revenue,
                "ebitda": ebitda,
                "gross_profit": gross_profit,
                "ev": ev,
                "market_cap": market_cap,
                "ev_rev": ev_rev,
                "ev_ebitda": ev_ebitda,
                "ev_gp": ev_gp,
                "pe_ratio": pe_ratio,
                "rev_growth": rev_growth,
                "gross_margin": gross_margin,
                "employees": info.get("fullTimeEmployees"),
            }
        )

    except Exception:
        return {}


def get_candidate_financials(candidate: Dict[str, Any]) -> Dict[str, Any]:
    ticker = candidate["ticker"].upper()

    # Prefer your curated internal financials database when available.
    fin = FINANCIALS_DB.get(ticker)

    if fin:
        return normalize_financial_record(fin)

    # Fallback lets the expanded universe work even if financials_db.py does not
    # contain every new ticker yet.
    return fetch_yfinance_financials(ticker, candidate.get("yf_symbol"))


def score_company(pub: Dict[str, Any], priv: Dict[str, Any], fin: Dict[str, Any]) -> float:
    relationship = pub.get("relationship", DIRECT)
    selected_sector = infer_sector(priv)

    score = 0.0

    if relationship == DIRECT:
        score += 45
    else:
        score += 22

    if pub.get("sector") == selected_sector:
        score += 20

    # Preserve intentional ordering within each sector.
    score += max(0, 10 - float(pub.get("order_index") or 0))

    priv_rev = safe_number(priv.get("revenue_m"), 0) * 1e6
    pub_rev = safe_number(fin.get("revenue"), 0)

    if priv_rev > 0 and pub_rev > 0:
        log_dist = abs(np.log10(pub_rev / priv_rev))
        score += max(0, 18 - log_dist * 7)

    priv_growth = safe_number(priv.get("rev_growth_pct"))
    pub_growth = safe_number(fin.get("rev_growth"))

    if priv_growth is not None and pub_growth is not None:
        priv_growth_decimal = priv_growth / 100
        diff = abs(pub_growth - priv_growth_decimal)
        score += max(0, 12 - diff * 60)

    priv_gm = safe_number(priv.get("gross_margin_pct"))
    pub_gm = safe_number(fin.get("gross_margin"))

    if priv_gm is not None and pub_gm is not None:
        priv_gm_decimal = priv_gm / 100
        diff = abs(pub_gm - priv_gm_decimal)
        score += max(0, 10 - diff * 40)

    if pub.get("status") == HISTORICAL:
        score -= 25

    if relationship == STRATEGIC:
        score = min(score, 72)

    return round(max(0, min(score, 99)), 1)



class PrivateCompanyInput(BaseModel):
    company_name: str
    sector: str
    sub_sector: Optional[str] = ""
    geo: Optional[str] = "North America"
    stage: Optional[str] = "Growth"

    revenue_m: float
    ebitda_m: Optional[float] = None
    gross_margin_pct: Optional[float] = None
    net_income_m: Optional[float] = None
    rev_growth_pct: Optional[float] = None
    employees: Optional[int] = None
    max_comps: Optional[int] = 8
    include_historical_comps: Optional[bool] = False

    # Extra banker/SaaS fields.
    business_description: Optional[str] = ""
    revenue_model: Optional[str] = ""
    customer_type: Optional[str] = ""
    arr_m: Optional[float] = None
    net_revenue_retention_pct: Optional[float] = None
    free_cash_flow_margin_pct: Optional[float] = None
    sales_efficiency: Optional[float] = None
    cac_payback_months: Optional[float] = None


def model_to_dict(model: BaseModel) -> Dict[str, Any]:
    if hasattr(model, "model_dump"):
        return model.model_dump()

    return model.dict()


@app.post("/api/find-comps")
def find_comps(inp: PrivateCompanyInput):
    priv = model_to_dict(inp)
    selected_sector = infer_sector(priv)
    include_historical = bool(priv.get("include_historical_comps", False))

    candidate_pool = build_candidate_pool(
        selected_sector=selected_sector,
        include_historical=include_historical,
    )

    scored = []

    for pub in candidate_pool:
        fin = get_candidate_financials(pub)

        # Skip names where we have neither curated DB data nor a yfinance fallback.
        if not fin:
            continue

        s = score_company(pub, priv, fin)

        universe_meta = get_company_universe_meta(pub["ticker"])

        scored.append(
            {
                **universe_meta,
                **pub,
                **fin,
                "match_score": s,
                "matchScore": s,
                "selectedSector": selected_sector,
            }
        )

    if not scored:
        return {
            "error": "No comparables found. Add these tickers to financials_db.py or make sure yfinance is available.",
            "selectedSector": selected_sector,
            "sector_label": selected_sector,
            "comps": [],
            "directComps": [],
            "strategicComps": [],
        }

    direct_scored = [c for c in scored if c.get("relationship") == DIRECT]
    strategic_scored = [c for c in scored if c.get("relationship") == STRATEGIC]

    direct_scored.sort(key=lambda x: x["match_score"], reverse=True)
    strategic_scored.sort(key=lambda x: x["match_score"], reverse=True)

    direct_limit = max(1, int(inp.max_comps or 8))
    strategic_limit = 5 if selected_sector == "AI / Machine Learning" else 0

    top_direct = direct_scored[:direct_limit]
    top_strategic = strategic_scored[:strategic_limit]
    top = top_direct + top_strategic

    if not top:
        return {"error": "No comparables found."}

    def stats(vals):
        vals = [safe_number(v) for v in vals]
        vals = [v for v in vals if v is not None and v > 0]

        if not vals:
            return None

        return {
            "median": round(float(np.median(vals)), 2),
            "p25": round(float(np.percentile(vals, 25)), 2),
            "p75": round(float(np.percentile(vals, 75)), 2),
            "mean": round(float(np.mean(vals)), 2),
            "min": round(float(np.min(vals)), 2),
            "max": round(float(np.max(vals)), 2),
        }

    # Valuation should be based on direct operating comps only.
    # Strategic ecosystem names like NVDA / TSM are shown for context but excluded
    # from the implied valuation range.
    valuation_base = [
        c
        for c in top_direct
        if c.get("relationship") == DIRECT and c.get("status") != HISTORICAL
    ]

    if not valuation_base:
        valuation_base = top_direct

    multiples = {
        "ev_rev": stats([c.get("ev_rev") for c in valuation_base]),
        "ev_ebitda": stats([c.get("ev_ebitda") for c in valuation_base]),
        "ev_gp": stats([c.get("ev_gp") for c in valuation_base]),
        "pe": stats([c.get("pe_ratio") for c in valuation_base]),
    }

    rev_m = inp.revenue_m
    ebitda_m = inp.ebitda_m
    gp_m = rev_m * (inp.gross_margin_pct / 100) if inp.gross_margin_pct else None

    implied = {}

    if multiples["ev_rev"] and rev_m:
        m = multiples["ev_rev"]
        implied["ev_rev"] = {
            "method": "EV / Revenue",
            "low": round(m["p25"] * rev_m, 1),
            "mid": round(m["median"] * rev_m, 1),
            "high": round(m["p75"] * rev_m, 1),
            "multiple_used": m["median"],
            "label": f"${rev_m}M x {m['median']}x",
        }

    if multiples["ev_ebitda"] and ebitda_m and ebitda_m > 0:
        m = multiples["ev_ebitda"]
        implied["ev_ebitda"] = {
            "method": "EV / EBITDA",
            "low": round(m["p25"] * ebitda_m, 1),
            "mid": round(m["median"] * ebitda_m, 1),
            "high": round(m["p75"] * ebitda_m, 1),
            "multiple_used": m["median"],
            "label": f"${ebitda_m}M x {m['median']}x",
        }

    if multiples["ev_gp"] and gp_m and gp_m > 0:
        m = multiples["ev_gp"]
        implied["ev_gp"] = {
            "method": "EV / Gross Profit",
            "low": round(m["p25"] * gp_m, 1),
            "mid": round(m["median"] * gp_m, 1),
            "high": round(m["p75"] * gp_m, 1),
            "multiple_used": m["median"],
            "label": f"${round(gp_m, 1)}M x {m['median']}x",
        }

    all_lows = [v["low"] for v in implied.values()]
    all_highs = [v["high"] for v in implied.values()]

    overall_range = (
        {
            "low": round(min(all_lows), 1),
            "high": round(max(all_highs), 1),
        }
        if all_lows
        else None
    )

    comps_out = []

    for c in top:
        revenue = safe_number(c.get("revenue"))
        ebitda = safe_number(c.get("ebitda"))
        ev = safe_number(c.get("ev"))
        market_cap = safe_number(c.get("market_cap"))

        revenue_m = round(revenue / 1e6, 1) if revenue else None
        ebitda_m_out = round(ebitda / 1e6, 1) if ebitda else None
        ev_b = round(ev / 1e9, 2) if ev else None
        market_cap_b = round(market_cap / 1e9, 2) if market_cap else None

        ebitda_margin = None
        if revenue and ebitda:
            ebitda_margin = round((ebitda / revenue) * 100, 1)

        rev_growth = (
            round(c["rev_growth"] * 100, 1)
            if c.get("rev_growth") is not None
            else None
        )

        gross_margin = (
            round(c["gross_margin"] * 100, 1)
            if c.get("gross_margin") is not None
            else None
        )

        relationship = c.get("relationship", DIRECT)
        status = c.get("status", LIVE)
        sector_label = c.get("sector") or selected_sector

        comps_out.append(
            {
                "ticker": c["ticker"],
                "symbol": c["ticker"],
                "name": c["name"],
                "companyName": c["name"],
                "sub": c.get("sub", sector_label),
                "industry": c.get("industry", sector_label),
                "sector": sector_label,
                "subSector": c.get("subSector", sector_label),
                "geography": c.get("geo", "Global"),
                "businessModel": (
                    "Strategic AI ecosystem comparable"
                    if relationship == STRATEGIC
                    else "Direct public company comparable"
                ),
                "description": c.get("description", c.get("reason", "")),
                "reason": c.get("reason", ""),
                "relationship": relationship,
                "compType": relationship,
                "status": status,
                "isStrategic": relationship == STRATEGIC,
                "isHistorical": status == HISTORICAL,
                "includedInValuation": relationship == DIRECT and status != HISTORICAL,
                "match_score": c["match_score"],
                "matchScore": c["match_score"],
                "market_cap_b": market_cap_b,
                "marketCap": market_cap,
                "ev_b": ev_b,
                "enterpriseValue": ev,
                "revenue_m": revenue_m,
                "revenue": revenue_m,
                "ebitda_m": ebitda_m_out,
                "ebitda": ebitda_m_out,
                "ev_rev": c.get("ev_rev"),
                "evRevenue": c.get("ev_rev"),
                "ev_ebitda": c.get("ev_ebitda"),
                "evEbitda": c.get("ev_ebitda"),
                "ev_gp": c.get("ev_gp"),
                "pe_ratio": c.get("pe_ratio"),
                "rev_growth": rev_growth,
                "revenueGrowth": rev_growth,
                "gross_margin": gross_margin,
                "grossMargin": gross_margin,
                "ebitda_margin": ebitda_margin,
                "ebitdaMargin": ebitda_margin,
                "employees": c.get("employees"),
                "marketDataDate": "Latest available market/database data",
                "financialPeriod": "Latest available fiscal year / LTM where available",
            }
        )

    direct_comps_out = [c for c in comps_out if c.get("relationship") == DIRECT]
    strategic_comps_out = [c for c in comps_out if c.get("relationship") == STRATEGIC]

    return {
        "comps": comps_out,
        "results": comps_out,
        "directComps": direct_comps_out,
        "strategicComps": strategic_comps_out,
        "multiples": multiples,
        "implied": implied,
        "overall_range": overall_range,
        "sector_label": selected_sector,
        "selectedSector": selected_sector,
        "comps_count": len(comps_out),
        "valuation_comps_count": len(valuation_base),
        "allSectors": SECTOR_LIST,
        "methodology": {
            "directComps": "Used for valuation multiples and implied valuation range.",
            "strategicComps": "Shown for ecosystem context only; excluded from direct valuation multiples.",
            "historicalComps": "Excluded by default unless include_historical_comps is true.",
        },
    }

@app.get("/api/health")
def health():
    return {"status": "ok", "companies_in_db": len(FINANCIALS_DB)}


@app.get("/api/sectors")
def sectors():
    return {"sectors": SECTOR_LIST}


FALLBACK_MA_DEALS = [
    {
        "buyer": "Salesforce",
        "target": "Informatica",
        "value": "$8.0B",
        "sector": "Data Infrastructure",
        "status": "Announced",
        "date": "2025-05-27",
        "angle": "Strategic data cloud expansion across AI, governance, and integration.",
        "sourceName": "Salesforce Newsroom",
        "sourceUrl": "https://www.salesforce.com/news/press-releases/2025/05/27/salesforce-signs-definitive-agreement-to-acquire-informatica/",
    },
    {
        "buyer": "HPE",
        "target": "Juniper Networks",
        "value": "$14.0B",
        "sector": "AI Networking",
        "status": "Closed",
        "date": "2025-07-02",
        "angle": "Networking scale-up for AI-native enterprise and cloud infrastructure.",
        "sourceName": "HPE Newsroom",
        "sourceUrl": "https://www.hpe.com/us/en/newsroom/press-release/2025/07/hewlett-packard-enterprise-completes-acquisition-of-juniper-networks.html",
    },
    {
        "buyer": "Synopsys",
        "target": "Ansys",
        "value": "$35.0B",
        "sector": "Engineering Software",
        "status": "Closed",
        "date": "2025-07-17",
        "angle": "Silicon-to-systems platform consolidation across EDA and simulation.",
        "sourceName": "Synopsys Newsroom",
        "sourceUrl": "https://news.synopsys.com/2025-07-17-Synopsys-Completes-Acquisition-of-Ansys",
    },
    {
        "buyer": "Cisco",
        "target": "Splunk",
        "value": "$28.0B",
        "sector": "Observability",
        "status": "Closed",
        "date": "2024-03-18",
        "angle": "Security, observability, and data analytics consolidation.",
        "sourceName": "Cisco Newsroom",
        "sourceUrl": "https://newsroom.cisco.com/c/r/newsroom/en/us/a/y2024/m03/cisco-completes-acquisition-of-splunk.html",
    },
    {
        "buyer": "IBM",
        "target": "HashiCorp",
        "value": "$6.4B",
        "sector": "Cloud Infrastructure",
        "status": "Closed",
        "date": "2025-02-27",
        "angle": "Hybrid cloud automation and infrastructure lifecycle management.",
        "sourceName": "IBM Newsroom",
        "sourceUrl": "https://newsroom.ibm.com/2025-02-27-IBM-Completes-Acquisition-of-HashiCorp",
    },
    {
        "buyer": "AMD",
        "target": "ZT Systems",
        "value": "$4.9B",
        "sector": "AI Infrastructure",
        "status": "Closed",
        "date": "2025-03-31",
        "angle": "AI systems design capabilities for hyperscale infrastructure.",
        "sourceName": "AMD Newsroom",
        "sourceUrl": "https://www.amd.com/en/newsroom/press-releases/2025-3-31-amd-completes-acquisition-of-zt-systems.html",
    },
    {
        "buyer": "ServiceNow",
        "target": "Moveworks",
        "value": "$2.85B",
        "sector": "Agentic AI",
        "status": "Announced",
        "date": "2025-03-10",
        "angle": "Enterprise AI assistant and automation platform expansion.",
        "sourceName": "ServiceNow Newsroom",
        "sourceUrl": "https://www.servicenow.com/company/media/press-room/servicenow-to-acquire-moveworks.html",
    },
    {
        "buyer": "Google",
        "target": "Wiz",
        "value": "$32.0B",
        "sector": "Cloud Security",
        "status": "Announced",
        "date": "2025-03-18",
        "angle": "Cloud-native security platform expansion for Google Cloud.",
        "sourceName": "Google Cloud Blog",
        "sourceUrl": "https://cloud.google.com/blog/products/identity-security/google-agreement-to-acquire-wiz",
    },
]


def normalize_fmp_deal(raw: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    if not isinstance(raw, dict):
        return None

    buyer = get_payload_field(raw, ["acquirerName", "acquirer", "buyer", "companyName"])
    target = get_payload_field(raw, ["targetName", "target", "name"])

    if not buyer or not target:
        return None

    transaction_value = get_payload_field(
        raw,
        ["transactionValue", "dealValue", "value", "totalTransactionValue"],
        "Undisclosed",
    )

    if isinstance(transaction_value, (int, float)):
        value = f"${transaction_value / 1_000_000_000:.1f}B"
    else:
        value = str(transaction_value)

    return {
        "buyer": buyer,
        "target": target,
        "value": value,
        "sector": get_payload_field(raw, ["sector", "industry"], "M&A"),
        "status": get_payload_field(raw, ["status"], "Latest"),
        "date": get_payload_field(raw, ["date", "announcedDate", "acceptedDate"], ""),
        "angle": get_payload_field(raw, ["description", "rationale"], "Latest M&A transaction from market data feed."),
        "sourceName": get_payload_field(raw, ["source", "sourceName"], "Market data source"),
        "sourceUrl": get_payload_field(raw, ["url", "sourceUrl", "link"], ""),
    }


@app.get("/api/market-intelligence/deals")
def latest_deals(limit: int = 6):
    fmp_api_key = os.getenv("FMP_API_KEY") or os.getenv("REACT_APP_FMP_API_KEY")

    if fmp_api_key:
        try:
            response = requests.get(
                "https://financialmodelingprep.com/stable/mergers-acquisitions-latest",
                params={"apikey": fmp_api_key, "limit": limit},
                timeout=8,
            )
            response.raise_for_status()
            raw_deals = response.json()
            deals = [normalize_fmp_deal(item) for item in raw_deals if isinstance(item, dict)]
            deals = [deal for deal in deals if deal]

            if deals:
                return {
                    "source": "fmp",
                    "updatedAt": datetime.now(timezone.utc).isoformat(),
                    "deals": deals[:limit],
                }
        except Exception:
            pass

    return {
        "source": "official-fallback",
        "updatedAt": datetime.now(timezone.utc).isoformat(),
        "deals": FALLBACK_MA_DEALS[: max(1, min(limit, len(FALLBACK_MA_DEALS)))],
    }


SEC_USER_AGENT = os.getenv(
    "SEC_USER_AGENT",
    "ValenceApp/1.0 veera.sachinsk@gmail.com",
)

SEC_HEADERS = {
    "User-Agent": SEC_USER_AGENT,
    "Accept-Encoding": "gzip, deflate",
}

TICKER_CIK_URL = "https://www.sec.gov/files/company_tickers.json"


@lru_cache(maxsize=1)
def load_ticker_to_cik_map():
    response = requests.get(TICKER_CIK_URL, headers=SEC_HEADERS, timeout=20)
    response.raise_for_status()

    data = response.json()
    ticker_map = {}

    for item in data.values():
        ticker = item["ticker"].upper()
        cik = str(item["cik_str"]).zfill(10)
        ticker_map[ticker] = cik

    return ticker_map


def sec_get_json(url):
    response = requests.get(url, headers=SEC_HEADERS, timeout=20)

    if response.status_code == 404:
        raise HTTPException(status_code=404, detail="SEC data not found.")

    response.raise_for_status()
    return response.json()


def build_filing_url(cik, accession_number, primary_document):
    cik_no_leading_zeros = str(int(cik))
    accession_no_dashes = accession_number.replace("-", "")

    return (
        f"https://www.sec.gov/Archives/edgar/data/"
        f"{cik_no_leading_zeros}/{accession_no_dashes}/{primary_document}"
    )


def get_recent_filings(submissions):
    recent = submissions.get("filings", {}).get("recent", {})

    forms = recent.get("form", [])
    filing_dates = recent.get("filingDate", [])
    report_dates = recent.get("reportDate", [])
    accession_numbers = recent.get("accessionNumber", [])
    primary_docs = recent.get("primaryDocument", [])

    important_forms = {
        "10-K",
        "10-Q",
        "8-K",
        "S-1",
        "DEF 14A",
        "10-K/A",
        "10-Q/A",
        "424B4",
    }

    filings = []

    for i, form in enumerate(forms):
        if form not in important_forms:
            continue

        accession = accession_numbers[i]
        primary_doc = primary_docs[i]

        filings.append(
            {
                "form": form,
                "filingDate": filing_dates[i],
                "reportDate": report_dates[i] if i < len(report_dates) else None,
                "accessionNumber": accession,
                "document": primary_doc,
                "url": build_filing_url(
                    submissions.get("cik"),
                    accession,
                    primary_doc,
                ),
            }
        )

        if len(filings) >= 25:
            break

    return filings


def latest_fact(companyfacts, possible_tags, unit="USD"):
    facts_root = companyfacts.get("facts", {}).get("us-gaap", {})

    for tag in possible_tags:
        concept = facts_root.get(tag)

        if not concept:
            continue

        unit_data = concept.get("units", {}).get(unit)

        if not unit_data:
            continue

        clean_items = [
            item
            for item in unit_data
            if item.get("val") is not None and item.get("filed")
        ]

        if not clean_items:
            continue

        latest = sorted(
            clean_items,
            key=lambda x: (x.get("filed", ""), x.get("end", "")),
            reverse=True,
        )[0]

        return {
            "tag": tag,
            "label": concept.get("label", tag),
            "value": latest.get("val"),
            "filed": latest.get("filed"),
            "period": latest.get("fp"),
            "form": latest.get("form"),
            "end": latest.get("end"),
        }

    return None


def get_financials(companyfacts):
    return {
        "revenue": latest_fact(
            companyfacts,
            [
                "RevenueFromContractWithCustomerExcludingAssessedTax",
                "SalesRevenueNet",
                "Revenues",
            ],
        ),
        "netIncome": latest_fact(
            companyfacts,
            [
                "NetIncomeLoss",
            ],
        ),
        "totalAssets": latest_fact(
            companyfacts,
            [
                "Assets",
            ],
        ),
        "totalLiabilities": latest_fact(
            companyfacts,
            [
                "Liabilities",
                "LiabilitiesCurrent",
            ],
        ),
        "stockholdersEquity": latest_fact(
            companyfacts,
            [
                "StockholdersEquity",
                "StockholdersEquityIncludingPortionAttributableToNoncontrollingInterest",
            ],
        ),
        "cash": latest_fact(
            companyfacts,
            [
                "CashAndCashEquivalentsAtCarryingValue",
                "CashCashEquivalentsRestrictedCashAndRestrictedCashEquivalents",
            ],
        ),
        "operatingIncome": latest_fact(
            companyfacts,
            [
                "OperatingIncomeLoss",
            ],
        ),
    }


def safe_float(value):
    try:
        if value is None:
            return None

        return float(value)
    except Exception:
        return None


def fast_info_get(fast_info, key):
    try:
        return fast_info[key]
    except Exception:
        try:
            return fast_info.get(key)
        except Exception:
            return None


def get_stock_data(ticker):
    try:
        stock = yf.Ticker(ticker)
        fast_info = stock.fast_info

        one_year_history = stock.history(period="1y")
        five_day_history = stock.history(period="5d")

        current_price = safe_float(fast_info_get(fast_info, "last_price"))

        if current_price is None and not five_day_history.empty:
            current_price = safe_float(five_day_history["Close"].dropna().iloc[-1])

        year_high = None
        year_low = None

        if not one_year_history.empty:
            year_high = safe_float(one_year_history["High"].max())
            year_low = safe_float(one_year_history["Low"].min())

        return {
            "ticker": ticker,
            "currentPrice": current_price,
            "previousClose": safe_float(fast_info_get(fast_info, "previous_close")),
            "dayHigh": safe_float(fast_info_get(fast_info, "day_high")),
            "dayLow": safe_float(fast_info_get(fast_info, "day_low")),
            "yearHigh": year_high,
            "yearLow": year_low,
            "marketCap": safe_float(fast_info_get(fast_info, "market_cap")),
            "enterpriseValue": safe_float(fast_info_get(fast_info, "enterprise_value")),
            "currency": fast_info_get(fast_info, "currency") or "USD",
        }

    except Exception as e:
        return {
            "ticker": ticker,
            "error": str(e),
        }


@app.get("/api/company/{ticker}")
def get_company_detail(ticker: str):
    ticker = ticker.upper().strip()

    ticker_map = load_ticker_to_cik_map()

    if ticker not in ticker_map:
        raise HTTPException(
            status_code=404,
            detail=f"No SEC CIK found for ticker {ticker}.",
        )

    cik = ticker_map[ticker]

    submissions = sec_get_json(
        f"https://data.sec.gov/submissions/CIK{cik}.json"
    )

    companyfacts = sec_get_json(
        f"https://data.sec.gov/api/xbrl/companyfacts/CIK{cik}.json"
    )

    return {
        "company": {
            "name": submissions.get("name"),
            "ticker": ticker,
            "cik": cik,
            "sic": submissions.get("sic"),
            "sicDescription": submissions.get("sicDescription"),
            "exchanges": submissions.get("exchanges", []),
            "entityType": submissions.get("entityType"),
        },
        "stock": get_stock_data(ticker),
        "financials": get_financials(companyfacts),
        "filings": get_recent_filings(submissions),
    }


def get_payload_field(obj: Dict[str, Any], keys: List[str], default=None):
    if not isinstance(obj, dict):
        return default

    for key in keys:
        value = obj.get(key)

        if value not in [None, "", "—", "N/A", "NA"]:
            return value

    return default


def prepare_private_company_for_report(raw_private: Dict[str, Any], selected_company: Dict[str, Any]):
    if not isinstance(raw_private, dict):
        raw_private = {}

    company_name = get_payload_field(
        raw_private,
        ["companyName", "company_name", "name", "targetName"],
        "Private Company",
    )

    sector = get_payload_field(
        raw_private,
        ["sector", "industry"],
        get_payload_field(selected_company, ["sector", "industry"], "N/A"),
    )

    sub_sector = get_payload_field(
        raw_private,
        ["subSector", "sub_sector", "sub", "industry"],
        sector,
    )

    geo = get_payload_field(
        raw_private,
        ["geo", "geography", "location", "headquarters"],
        "North America",
    )

    revenue_m = get_payload_field(raw_private, ["revenue_m", "revenue", "annualRevenue"])
    ebitda_m = get_payload_field(raw_private, ["ebitda_m", "ebitda", "EBITDA"])
    gross_margin = get_payload_field(raw_private, ["gross_margin_pct", "grossMargin", "gross_margin"])
    rev_growth = get_payload_field(raw_private, ["rev_growth_pct", "revenueGrowth", "revenue_growth"])
    employees = get_payload_field(raw_private, ["employees", "employeeCount", "employee_count"])

    fcf_margin = get_payload_field(
        raw_private,
        ["free_cash_flow_margin_pct", "fcfMargin", "freeCashFlowMargin"],
    )

    rule_of_40 = None
    rev_growth_num = safe_number(rev_growth)
    fcf_margin_num = safe_number(fcf_margin)

    if rev_growth_num is not None and fcf_margin_num is not None:
        rule_of_40 = rev_growth_num + fcf_margin_num

    return {
        "companyName": company_name,
        "name": company_name,
        "description": get_payload_field(
            raw_private,
            ["business_description", "businessDescription", "description"],
            "Not provided",
        ),
        "businessDescription": get_payload_field(
            raw_private,
            ["business_description", "businessDescription", "description"],
            "Not provided",
        ),
        "sector": sector,
        "subSector": sub_sector,
        "industry": sub_sector,
        "geography": geo,
        "revenueModel": get_payload_field(
            raw_private,
            ["revenue_model", "revenueModel", "businessModel", "business_model"],
            "Not provided",
        ),
        "businessModel": get_payload_field(
            raw_private,
            ["revenue_model", "revenueModel", "businessModel", "business_model"],
            "Not provided",
        ),
        "customerType": get_payload_field(
            raw_private,
            ["customer_type", "customerType"],
            "Not provided",
        ),
        "revenue_m": revenue_m,
        "revenue": revenue_m,
        "ebitda_m": ebitda_m,
        "ebitda": ebitda_m,
        "grossMargin": gross_margin,
        "gross_margin": gross_margin,
        "revenueGrowth": rev_growth,
        "revenue_growth": rev_growth,
        "employees": employees,
        "arr": get_payload_field(raw_private, ["arr_m", "arr", "ARR"]),
        "netRevenueRetention": get_payload_field(
            raw_private,
            ["net_revenue_retention_pct", "netRevenueRetention", "nrr"],
        ),
        "freeCashFlowMargin": fcf_margin,
        "fcfMargin": fcf_margin,
        "ruleOf40": get_payload_field(raw_private, ["ruleOf40", "rule_of_40"], rule_of_40),
        "salesEfficiency": get_payload_field(
            raw_private,
            ["sales_efficiency", "salesEfficiency"],
        ),
        "cacPaybackMonths": get_payload_field(
            raw_private,
            ["cac_payback_months", "cacPaybackMonths"],
        ),
    }


def prepare_company_for_report(raw_company: Dict[str, Any]):
    if not isinstance(raw_company, dict):
        raw_company = {}

    ticker = get_payload_field(raw_company, ["ticker", "symbol", "Ticker"], "N/A")
    name = get_payload_field(raw_company, ["name", "companyName", "company", "Company"], ticker)

    revenue_m = get_payload_field(raw_company, ["revenue_m", "revenue", "Revenue"])
    ev_rev = get_payload_field(raw_company, ["ev_rev", "evRevenue", "ev_revenue"])
    ev_ebitda = get_payload_field(raw_company, ["ev_ebitda", "evEbitda", "ev_ebitda"])
    ev_b = get_payload_field(raw_company, ["ev_b"])
    market_cap_b = get_payload_field(raw_company, ["market_cap_b"])

    enterprise_value = get_payload_field(
        raw_company,
        ["enterpriseValue", "enterprise_value", "ev", "EV"],
    )

    if enterprise_value is None and ev_b is not None:
        enterprise_value = safe_number(ev_b) * 1_000_000_000

    market_cap = get_payload_field(raw_company, ["marketCap", "market_cap"])

    if market_cap is None and market_cap_b is not None:
        market_cap = safe_number(market_cap_b) * 1_000_000_000

    ebitda_m = get_payload_field(raw_company, ["ebitda_m", "ebitda", "EBITDA"])

    ebitda_margin = get_payload_field(
        raw_company,
        ["ebitdaMargin", "ebitda_margin"],
    )

    if ebitda_margin is None:
        ebitda_margin = safe_divide(ebitda_m, revenue_m)
        if ebitda_margin is not None:
            ebitda_margin *= 100

    gross_margin = get_payload_field(
        raw_company,
        ["grossMargin", "gross_margin", "gross_margin_pct"],
    )

    revenue_growth = get_payload_field(
        raw_company,
        ["revenueGrowth", "rev_growth", "revenue_growth", "rev_growth_pct"],
    )

    match_score = get_payload_field(
        raw_company,
        ["matchScore", "match_score", "score"],
        100,
    )

    sector = get_payload_field(raw_company, ["sector"], "N/A")
    industry = get_payload_field(raw_company, ["industry", "sub", "subSector"], sector)

    return {
        **raw_company,
        "ticker": ticker,
        "symbol": ticker,
        "name": name,
        "companyName": name,
        "sector": sector,
        "industry": industry,
        "subSector": get_payload_field(raw_company, ["subSector", "sub_sector", "sub"], industry),
        "description": get_payload_field(raw_company, ["description"], ""),
        "businessModel": get_payload_field(raw_company, ["businessModel", "business_model"], "Public company comparable"),
        "revenue_m": revenue_m,
        "revenue": revenue_m,
        "ebitda_m": ebitda_m,
        "ebitda": ebitda_m,
        "enterpriseValue": enterprise_value,
        "marketCap": market_cap,
        "ev_rev": ev_rev,
        "evRevenue": ev_rev,
        "ev_ebitda": ev_ebitda,
        "evEbitda": ev_ebitda,
        "ebitdaMargin": ebitda_margin,
        "ebitda_margin": ebitda_margin,
        "grossMargin": gross_margin,
        "gross_margin": gross_margin,
        "revenueGrowth": revenue_growth,
        "revenue_growth": revenue_growth,
        "matchScore": match_score,
        "match_score": match_score,
        "marketDataDate": get_payload_field(
            raw_company,
            ["marketDataDate", "asOfDate", "sourceDate"],
            "Latest available market/database data",
        ),
        "financialPeriod": get_payload_field(
            raw_company,
            ["financialPeriod", "period", "fiscalYear"],
            "Latest available fiscal year / LTM where available",
        ),
    }


def build_report_response(payload: Dict[str, Any]):
    selected_raw = (
        payload.get("selectedCompany")
        or payload.get("selectedComparable")
        or payload.get("company")
        or {}
    )

    raw_comps = payload.get("comps") or [selected_raw]

    if not isinstance(raw_comps, list):
        raw_comps = [selected_raw]

    selected_company = prepare_company_for_report(selected_raw)
    comps = [prepare_company_for_report(comp) for comp in raw_comps]

    selected_ticker = selected_company.get("ticker") or selected_company.get("symbol")

    if selected_ticker:
        existing_tickers = {str(comp.get("ticker") or comp.get("symbol")) for comp in comps}
        if str(selected_ticker) not in existing_tickers:
            comps = [selected_company] + comps

    private_company = prepare_private_company_for_report(
        payload.get("privateCompany") or {},
        selected_company,
    )

    clean_ticker = safe_report_name(selected_company.get("ticker"))
    filename = f"valence-report-{clean_ticker}.pdf"
    output_path = os.path.join(tempfile.gettempdir(), filename)

    try:
        generate_banker_report(
            selected_company=selected_company,
            comps=comps,
            private_company=private_company,
            output_path=output_path,
        )
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"Could not generate report: {str(e)}",
        )

    return FileResponse(
        output_path,
        media_type="application/pdf",
        filename=filename,
    )


@app.post("/api/generate-report")
async def generate_report(payload: Dict[str, Any]):
    return build_report_response(payload)


@app.post("/api/generate-deck")
async def generate_deck(payload: Dict[str, Any]):
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise HTTPException(
            status_code=500,
            detail="PowerPoint export is not installed on the backend.",
        ) from exc

    selected = payload.get("selectedCompany") or {}
    comps = payload.get("comps") or []
    private = payload.get("privateCompany") or {}
    implied = payload.get("implied") or {}
    overall = payload.get("overall_range") or {}
    sector = payload.get("sector_label") or selected.get("sector") or "selected sector"
    options = payload.get("deckOptions") or {}

    colorways = {
        "midnight": ("07111F", "0F1B2D", "F8FAFC", "94A3B8", "22D3EE", "7C3AED"),
        "boardroom": ("F6F8FB", "FFFFFF", "102033", "5D6B7A", "1D4ED8", "0F766E"),
        "emerald": ("071A16", "0D2B24", "ECFDF5", "A7F3D0", "10B981", "F59E0B"),
        "plum": ("170923", "251137", "FAF5FF", "D8B4FE", "A855F7", "22D3EE"),
        "onyx": ("08070D", "15111F", "FAFAFA", "A3A3A3", "F472B6", "38BDF8"),
        "copper": ("12100D", "211A14", "FFF7ED", "FED7AA", "F97316", "14B8A6"),
        "arctic": ("EFF6FF", "FFFFFF", "0F172A", "475569", "0284C7", "7C3AED"),
        "terminal": ("03120C", "082016", "ECFDF5", "86EFAC", "22C55E", "38BDF8"),
    }

    bg, panel, ink, muted, accent, accent2 = colorways.get(
        options.get("colorway"),
        colorways["midnight"],
    )

    def rgb(hex_value: str):
        clean = str(hex_value).replace("#", "")
        return RGBColor(int(clean[0:2], 16), int(clean[2:4], 16), int(clean[4:6], 16))

    def text_value(obj: Dict[str, Any], keys: List[str], fallback: str = "N/A") -> str:
        for key in keys:
            value = obj.get(key)
            if value not in (None, ""):
                return str(value)
        return fallback

    def number_value(obj: Dict[str, Any], keys: List[str]):
        for key in keys:
            try:
                value = obj.get(key)
                if value not in (None, ""):
                    return float(value)
            except (TypeError, ValueError):
                continue
        return None

    def fmt_m(value):
        try:
            return f"${float(value):,.0f}M"
        except (TypeError, ValueError):
            return "N/A"

    def fmt_x(value):
        try:
            return f"{float(value):.1f}x"
        except (TypeError, ValueError):
            return "N/A"

    def fmt_pct(value):
        try:
            return f"{float(value):.0f}%"
        except (TypeError, ValueError):
            return "N/A"

    def add_text(slide, text, x, y, w, h, size=18, bold=False, color=None, align=None):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        paragraph = frame.paragraphs[0]
        if align:
            paragraph.alignment = align
        run = paragraph.add_run()
        run.text = str(text)
        run.font.name = "Aptos"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color or ink)
        return box

    def add_panel(slide, x, y, w, h, fill=None):
        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill or panel)
        shape.line.color.rgb = rgb(accent)
        shape.line.transparency = 35
        return shape

    def add_metric_card(slide, label, value, note, x, y, w=2.75, h=1.18):
        add_panel(slide, x, y, w, h)
        add_text(slide, label, x + 0.18, y + 0.16, w - 0.35, 0.22, 9, True, accent)
        add_text(slide, value, x + 0.18, y + 0.44, w - 0.35, 0.36, 20, True, ink)
        add_text(slide, note, x + 0.18, y + 0.84, w - 0.35, 0.2, 8, False, muted)

    def add_bullet_rows(slide, rows, start_y=1.55, max_rows=5):
        for index, item in enumerate(rows[:max_rows]):
            y = start_y + index * 0.86
            add_panel(slide, 0.78, y, 11.75, 0.58)
            add_text(slide, item, 1.04, y + 0.16, 11.1, 0.22, 13, index == 0, ink)

    def add_bar(slide, label, value, max_value, x, y, w=5.3):
        safe_value = max(0, float(value or 0))
        safe_max = max(float(max_value or 1), safe_value, 1)
        bar_w = max(0.08, w * safe_value / safe_max)
        add_text(slide, label, x, y, 2.2, 0.22, 9, True, muted)
        add_panel(slide, x + 2.35, y + 0.02, w, 0.16, "1E293B")
        add_panel(slide, x + 2.35, y + 0.02, bar_w, 0.16, accent)
        add_text(slide, f"{safe_value:.1f}x", x + 2.35 + w + 0.15, y - 0.03, 0.7, 0.22, 9, True, ink)

    def add_quadrant(slide, x, y, w, h, points):
        add_panel(slide, x, y, w, h)
        add_text(slide, "Lower fit", x + 0.15, y + h - 0.28, 1.1, 0.18, 8, True, muted)
        add_text(slide, "Higher fit", x + w - 1.2, y + 0.12, 1.0, 0.18, 8, True, muted, PP_ALIGN.RIGHT)
        add_text(slide, "Valuation support", x + w - 1.55, y + h - 0.28, 1.35, 0.18, 8, True, muted, PP_ALIGN.RIGHT)
        add_panel(slide, x + w / 2, y + 0.16, 0.02, h - 0.32, muted)
        add_panel(slide, x + 0.16, y + h / 2, w - 0.32, 0.02, muted)
        for point in points[:7]:
            px = x + 0.35 + (w - 0.7) * max(0, min(100, point["fit"])) / 100
            py = y + h - 0.35 - (h - 0.7) * max(0, min(100, point["support"])) / 100
            dot = slide.shapes.add_shape(1, Inches(px), Inches(py), Inches(0.16), Inches(0.16))
            dot.fill.solid()
            dot.fill.fore_color.rgb = rgb(accent2 if point.get("lead") else accent)
            dot.line.color.rgb = rgb(accent2 if point.get("lead") else accent)
            add_text(slide, point["label"], px + 0.18, py - 0.02, 0.9, 0.18, 7, True, ink)

    def add_waterfall(slide, x, y, values):
        max_abs = max([abs(v["value"]) for v in values] + [1])
        base_y = y + 2.15
        cursor_x = x
        for item in values:
            height = 1.8 * abs(item["value"]) / max_abs
            top_y = base_y - height if item["value"] >= 0 else base_y
            fill = accent if item["value"] >= 0 else accent2
            add_panel(slide, cursor_x, top_y, 1.25, max(0.12, height), fill)
            add_text(slide, item["label"], cursor_x - 0.08, base_y + 0.18, 1.4, 0.34, 8, True, muted, PP_ALIGN.CENTER)
            add_text(slide, item["display"], cursor_x - 0.05, top_y - 0.26, 1.35, 0.2, 9, True, ink, PP_ALIGN.CENTER)
            cursor_x += 1.55

    def add_roadmap(slide, x, y, steps):
        for index, step in enumerate(steps):
            card_x = x + index * 3.05
            add_panel(slide, card_x, y, 2.72, 2.1)
            add_text(slide, f"0{index + 1}", card_x + 0.2, y + 0.18, 0.45, 0.22, 10, True, accent)
            add_text(slide, step["title"], card_x + 0.2, y + 0.5, 2.2, 0.38, 14, True, ink)
            add_text(slide, step["body"], card_x + 0.2, y + 1.04, 2.24, 0.62, 9, False, muted)
            if index < len(steps) - 1:
                add_text(slide, ">", card_x + 2.83, y + 0.84, 0.18, 0.24, 12, True, accent)

    def add_risk_table(slide, rows):
        headers = ["Risk", "Why it matters", "Mitigation"]
        col_x = [0.82, 3.55, 7.3]
        widths = [2.35, 3.2, 4.85]
        for x, width, header in zip(col_x, widths, headers):
            add_text(slide, header, x, 1.48, width, 0.24, 10, True, accent)
        for index, row in enumerate(rows[:5]):
            y = 1.92 + index * 0.78
            add_panel(slide, 0.72, y - 0.08, 11.65, 0.58)
            for x, width, value in zip(col_x, widths, row):
                add_text(slide, value, x, y + 0.06, width, 0.32, 9, False, ink)

    def add_slide(title, kicker="VALENCE EXPORT STUDIO"):
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = rgb(bg)
        add_text(slide, kicker, 0.7, 0.42, 4.8, 0.25, 10, True, accent)
        add_text(slide, title, 0.7, 0.72, 7.8, 0.55, 24, True, ink)
        return slide

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    target_name = text_value(private, ["company_name", "name", "company"], "Target company")
    selected_name = text_value(selected, ["name", "company", "companyName"], "Selected comp")
    selected_ticker = text_value(selected, ["ticker", "symbol"], "")
    audience = str(options.get("audience") or "investor").replace("-", " ")
    feature_pack = str(options.get("featurePack") or "competition").replace("-", " ")
    density = str(options.get("density") or "balanced")
    component_style = str(options.get("componentStyle") or "strategy").replace("-", " ")
    instructions = str(options.get("instructions") or "").strip()
    selected_sections = options.get("sections") or {}
    selected_sections = {
        "overview": True,
        "comps": True,
        "multiples": True,
        "valuation": True,
        "market": False,
        **selected_sections,
    }
    revenue_m = number_value(private, ["revenue_m", "revenue"])
    ebitda_m = number_value(private, ["ebitda_m", "ebitda"])
    growth_pct = number_value(private, ["rev_growth_pct", "revGrowth", "revenue_growth"])
    gross_margin = number_value(private, ["gross_margin_pct", "grossMargin", "gross_margin"])
    ev_rev_values = [
        number_value(comp, ["ev_rev", "evRevenue"])
        for comp in comps
        if number_value(comp, ["ev_rev", "evRevenue"]) is not None
    ]
    ev_ebitda_values = [
        number_value(comp, ["ev_ebitda", "evEbitda"])
        for comp in comps
        if number_value(comp, ["ev_ebitda", "evEbitda"]) is not None
    ]
    median_ev_rev = float(np.median(ev_rev_values)) if ev_rev_values else None
    median_ev_ebitda = float(np.median(ev_ebitda_values)) if ev_ebitda_values else None

    slide = add_slide(f"{target_name} Trading Comps Analysis", "CONFIDENTIAL DRAFT")
    add_text(slide, f"Primary benchmark: {selected_name} {f'({selected_ticker})' if selected_ticker else ''}", 0.7, 1.55, 8.2, 0.38, 17, False, muted)
    add_panel(slide, 8.7, 1.2, 3.7, 4.7)
    add_text(slide, "Deck setup", 9.05, 1.55, 2.7, 0.3, 11, True, accent)
    add_text(slide, f"Audience\n{audience.title()}", 9.05, 2.02, 2.7, 0.72, 18, True, ink)
    add_text(slide, f"Feature pack\n{feature_pack.title()}", 9.05, 2.95, 2.7, 0.72, 18, True, ink)
    add_text(slide, f"Style\n{component_style.title()} / {density.title()}", 9.05, 3.9, 2.7, 0.72, 18, True, ink)
    add_text(slide, f"Sector: {sector}", 0.7, 6.45, 6.2, 0.28, 12, True, accent)

    slide = add_slide("Executive valuation thesis")
    thesis = [
        f"{target_name} screens against {len(comps) or 'selected'} public comps in {sector}.",
        f"{selected_name} is the lead benchmark for business model and trading context.",
        "Valuation range is based on selected peer multiples and user-entered target financials.",
        f"Output is tailored for {audience} use with {feature_pack} support.",
    ]
    for index, item in enumerate(thesis):
        add_panel(slide, 0.8, 1.45 + index * 1.02, 11.7, 0.72)
        add_text(slide, item, 1.05, 1.63 + index * 1.02, 10.9, 0.3, 16, index == 0, ink)

    if selected_sections.get("overview"):
        slide = add_slide("Situation, question, and valuation answer")
        add_metric_card(slide, "Revenue", fmt_m(revenue_m), "User-entered target metric", 0.82, 1.45)
        add_metric_card(slide, "Growth", fmt_pct(growth_pct), "Revenue growth context", 3.85, 1.45)
        add_metric_card(slide, "Gross margin", fmt_pct(gross_margin), "Quality / margin signal", 6.88, 1.45)
        add_metric_card(slide, "Lead comp", selected_ticker or selected_name, "Primary public benchmark", 9.91, 1.45)
        situation_rows = [
            f"Core question: what public trading context best supports {target_name}'s valuation range?",
            f"Answer: use the {sector} peer set, anchored on {selected_name}, then pressure-test with growth and margin quality.",
            "Investor-ready narrative: valuation is most defensible when peer relevance, multiple dispersion, and business quality are shown together.",
            f"Deck mode: {feature_pack.title()} for {audience.title()} audiences.",
        ]
        add_bullet_rows(slide, situation_rows, 3.15, 4)

    if selected_sections.get("comps"):
        slide = add_slide("Comparable universe and selection logic")
        headers = ["Ticker", "Company", "Revenue", "EV/Rev", "EV/EBITDA", "Match"]
        rows = comps[:8] or [selected]
        col_x = [0.75, 1.75, 5.45, 7.15, 8.75, 10.65]
        widths = [0.8, 3.3, 1.3, 1.25, 1.5, 1.0]
        for x, width, header in zip(col_x, widths, headers):
            add_text(slide, header, x, 1.45, width, 0.28, 10, True, accent)
        for row_index, comp_row in enumerate(rows):
            y = 1.9 + row_index * 0.48
            add_panel(slide, 0.65, y - 0.08, 11.55, 0.38, panel)
            values = [
                text_value(comp_row, ["ticker", "symbol"], "N/A"),
                text_value(comp_row, ["name", "company", "companyName"], "N/A"),
                fmt_m(number_value(comp_row, ["revenue_m", "revenue"])),
                fmt_x(number_value(comp_row, ["ev_rev", "evRevenue"])),
                fmt_x(number_value(comp_row, ["ev_ebitda", "evEbitda"])),
                fmt_pct(number_value(comp_row, ["match_score", "matchScore"])),
            ]
            for x, width, value in zip(col_x, widths, values):
                add_text(slide, value, x, y, width, 0.22, 10, False, ink)
        add_text(slide, "Selection rule: prioritize business-model relevance, sector fit, valuation data completeness, and match score.", 0.8, 6.35, 11.4, 0.28, 11, True, muted)

    if selected_sections.get("multiples"):
        slide = add_slide("Trading multiple benchmark readout")
        add_metric_card(slide, "Median EV/Revenue", fmt_x(median_ev_rev), "Selected peer median", 0.82, 1.42, 3.2)
        add_metric_card(slide, "Median EV/EBITDA", fmt_x(median_ev_ebitda), "Selected peer median", 4.08, 1.42, 3.2)
        add_metric_card(slide, "Peer count", str(len(comps) or 1), "Companies in selected set", 7.34, 1.42, 3.2)
        multiple_rows = [
            "Premium multiple support requires more than category heat: growth, margin, scale, and public-market relevance must align.",
            f"{selected_name} anchors the discussion, while the broader comp set establishes the defensible range.",
            "Outliers should be explained rather than blindly averaged; appendix pages preserve the full evidence trail.",
        ]
        add_bullet_rows(slide, multiple_rows, 3.15, 3)

        slide = add_slide("Peer multiple benchmarking")
        ranked_ev_rev = sorted(
            [
                {
                    "label": text_value(comp, ["ticker", "symbol"], text_value(comp, ["name", "company"], "Comp"))[:10],
                    "value": number_value(comp, ["ev_rev", "evRevenue"]) or 0,
                }
                for comp in comps[:8]
            ],
            key=lambda item: item["value"],
            reverse=True,
        )
        max_ev_rev = max([item["value"] for item in ranked_ev_rev] + [1])
        for index, item in enumerate(ranked_ev_rev[:8]):
            add_bar(slide, item["label"], item["value"], max_ev_rev, 0.9, 1.45 + index * 0.46, 6.1)
        add_panel(slide, 9.0, 1.48, 2.65, 3.0)
        add_text(slide, "Readout", 9.25, 1.82, 1.8, 0.25, 12, True, accent)
        add_text(slide, "This page makes the multiple dispersion visible so the chosen range feels defended, not arbitrary.", 9.25, 2.28, 1.95, 0.95, 13, False, ink)

        slide = add_slide("Strategic fit vs. valuation support matrix")
        matrix_points = []
        for comp in comps[:7]:
            match = number_value(comp, ["match_score", "matchScore"]) or 55
            ev_rev = number_value(comp, ["ev_rev", "evRevenue"]) or median_ev_rev or 5
            support = 50 if not median_ev_rev else min(96, max(18, 50 + (ev_rev - median_ev_rev) * 7))
            label = text_value(comp, ["ticker", "symbol"], text_value(comp, ["name", "company"], "Comp"))[:8]
            matrix_points.append({
                "label": label,
                "fit": match,
                "support": support,
                "lead": label == selected_ticker,
            })
        add_quadrant(slide, 1.0, 1.45, 6.2, 4.45, matrix_points)
        add_panel(slide, 8.15, 1.65, 3.35, 3.7)
        add_text(slide, "How to use it", 8.45, 1.98, 2.4, 0.25, 12, True, accent)
        add_text(slide, "Comps in the upper-right are stronger valuation anchors. Lower-left comps may still be useful as boundary cases or appendix support.", 8.45, 2.42, 2.55, 1.2, 13, False, ink)

    if selected_sections.get("valuation"):
        slide = add_slide("Implied valuation range and method bridge")
        low = overall.get("low")
        high = overall.get("high")
        if low is not None and high is not None:
            add_panel(slide, 0.85, 1.55, 5.2, 2.2)
            add_text(slide, "Enterprise value range", 1.15, 1.9, 3.2, 0.28, 12, True, accent)
            add_text(slide, f"{fmt_m(low)} - {fmt_m(high)}", 1.15, 2.38, 4.3, 0.55, 30, True, ink)
        else:
            add_text(slide, "Add revenue, EBITDA, and margin inputs to populate the implied valuation range.", 0.85, 1.7, 9.5, 0.45, 18, True, muted)
        methods = list(implied.values())[:4]
        for index, method in enumerate(methods):
            y = 4.25 + index * 0.5
            label = method.get("method") or "Valuation method"
            add_text(slide, f"{label}: {fmt_m(method.get('low'))} - {fmt_m(method.get('high'))}", 0.95, y, 7.6, 0.25, 12, False, ink)
        add_panel(slide, 7.1, 1.55, 4.95, 2.2)
        add_text(slide, "Valuation defense", 7.4, 1.9, 3.7, 0.28, 12, True, accent)
        add_text(slide, "Triangulate revenue, EBITDA, and gross-profit methods where available; explain gaps where target metrics are unavailable.", 7.4, 2.35, 4.1, 0.7, 14, False, ink)

        slide = add_slide("Valuation bridge from target metric to EV")
        bridge_values = [
            {"label": "Revenue", "value": revenue_m or 0, "display": fmt_m(revenue_m)},
            {"label": "Peer multiple", "value": median_ev_rev or 0, "display": fmt_x(median_ev_rev)},
            {"label": "Low case", "value": low or 0, "display": fmt_m(low)},
            {"label": "High case", "value": high or 0, "display": fmt_m(high)},
        ]
        add_waterfall(slide, 1.0, 2.0, bridge_values)
        add_panel(slide, 8.2, 1.65, 3.4, 3.15)
        add_text(slide, "Bridge logic", 8.5, 1.98, 2.1, 0.25, 12, True, accent)
        add_text(slide, "The bridge turns the valuation from a number into a defendable chain: target financial metric, peer multiple, and sensitivity range.", 8.5, 2.42, 2.55, 1.05, 13, False, ink)

    slide = add_slide("Recommendation and executive talking points")
    recommendation_rows = [
        f"Recommendation: position {target_name} against the selected {sector} peer set, with {selected_name} as the lead public benchmark.",
        "Use median multiples for the base case and 25th-75th percentile ranges for sensitivity discussion.",
        "Lead with business-model fit before valuation math; decision makers trust the multiple only after they trust the comp set.",
        "Reserve outlier discussion, source limitations, and alternate comps for appendix support.",
    ]
    add_bullet_rows(slide, recommendation_rows, 1.45, 4)

    slide = add_slide("Likely Q&A and objection handling")
    qa_rows = [
        "Why these comps? Selection emphasizes sector, business model, revenue scale, growth, margin profile, and public-market data quality.",
        "Why not a higher multiple? Premium support depends on proving growth durability, margin expansion, and strategic scarcity.",
        "What weakens the valuation? Thin profitability, volatile growth, missing retention data, or weak comp relevance.",
        "What would strengthen the case? More KPIs, customer concentration detail, ARR/NRR, cohort retention, and a clearer market-size bridge.",
    ]
    add_bullet_rows(slide, qa_rows, 1.45, 4)

    slide = add_slide("Risk register and mitigation plan")
    add_risk_table(slide, [
        ["Comp relevance", "Wrong peer set weakens the multiple", "Keep alternates in appendix and explain inclusion rules"],
        ["Profitability gap", "High growth with low margin can pressure valuation", "Show margin path and operating leverage cases"],
        ["Growth durability", "Multiple depends on credible forward growth", "Add ARR, NRR, pipeline, and cohort evidence"],
        ["Market volatility", "Public multiples can reset quickly", "Use percentile ranges and update data date"],
        ["Data gaps", "Missing KPIs reduce confidence", "Flag assumptions and request diligence inputs"],
    ])

    slide = add_slide("Workplan to final investment-ready deck")
    add_roadmap(slide, 0.9, 1.75, [
        {"title": "Validate", "body": "Confirm company profile, revenue model, sector, and excluded comps."},
        {"title": "Pressure-test", "body": "Run multiple, growth, and margin sensitivities against selected peers."},
        {"title": "Narrate", "body": "Turn valuation output into recommendation, risks, and buyer/investor story."},
        {"title": "Finalize", "body": "Export PPTX, PDF, or Google Slides-ready file with appendix support."},
    ])

    if selected_sections.get("market") or feature_pack in {"market", "consulting"}:
        slide = add_slide("Market intelligence appendix")
        market_rows = [
            f"Sector signal: {sector} valuation context should be interpreted alongside buyer activity and public comp sentiment.",
            "Strategic M&A examples help explain why certain categories command higher multiples.",
            "Use this appendix to support the 'why now' narrative without overloading the executive recommendation.",
        ]
        add_bullet_rows(slide, market_rows, 1.45, 3)

    slide = add_slide("Appendix: data quality and model notes")
    appendix_rows = [
        "Financial data uses latest available public market and company database values available to Valence.",
        "Implied valuation ranges are directional and should be paired with diligence on ARR, NRR, margin profile, and customer quality.",
        "Exported slides are generated from user-selected comps; changing selected companies changes the valuation evidence base.",
        "This deck is original Valence output and does not copy third-party consulting presentation templates.",
    ]
    add_bullet_rows(slide, appendix_rows, 1.45, 4)

    slide = add_slide("Recommended next steps")
    next_steps = [
        "Validate selected comps against product scope, buyer, and end-market exposure.",
        "Pressure-test valuation range with growth, profitability, and retention sensitivities.",
        "Prepare appendix support for top objections around multiple selection.",
    ]
    if instructions:
        next_steps.append(f"Custom brief: {instructions[:150]}")
    for index, item in enumerate(next_steps[:4]):
        add_panel(slide, 0.85, 1.45 + index * 0.95, 11.5, 0.66)
        add_text(slide, item, 1.1, 1.62 + index * 0.95, 10.8, 0.28, 15, index == 0, ink)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    raw_name = text_value(private, ["company_name", "name", "company"], selected_ticker or "company")
    filename = re.sub(r"[^A-Za-z0-9_-]+", "-", raw_name).strip("-").lower() or "company"

    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="valence-deck-{filename}.pptx"'},
    )


@app.post("/api/download-report")
async def download_report(payload: Dict[str, Any]):
    return build_report_response(payload)
