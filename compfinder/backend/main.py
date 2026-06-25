from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.responses import FileResponse, StreamingResponse
from pydantic import BaseModel
from typing import Optional, Any, Dict, List
import numpy as np
import os
import requests
import yfinance as yf
from functools import lru_cache
from datetime import datetime
import io
import sys
import tempfile
from pathlib import Path

from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.units import inch

from company_universe import COMPANY_UNIVERSE
from financials_db import FINANCIALS_DB

ROOT_BACKEND = Path(__file__).resolve().parents[2] / "backend"
if str(ROOT_BACKEND) not in sys.path:
    sys.path.insert(0, str(ROOT_BACKEND))

from report_generator import generate_banker_report

app = FastAPI(title="CompFinder API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

INPUT_SECTOR_MAP = {
    "Enterprise SaaS / Tech":      "saas_tech",
    "Healthcare / Biotech":        "healthcare",
    "Financial Services":          "financial",
    "Consumer / Retail":           "consumer",
    "Industrials / Manufacturing": "industrials",
}

SECTOR_LABEL_MAP = {v: k for k, v in INPUT_SECTOR_MAP.items()}

SEC_USER_AGENT = os.getenv("SEC_USER_AGENT", "ValenceApp/1.0 veera.sachinsk@gmail.com")
SEC_HEADERS = {
    "User-Agent": SEC_USER_AGENT,
    "Accept-Encoding": "gzip, deflate",
}
TICKER_CIK_URL = "https://www.sec.gov/files/company_tickers.json"
MA_UNIVERSE_LIMIT = 10000


# ── Scoring ───────────────────────────────────────────────────────────────────
def score_company(pub, priv, fin):
    score = 0.0
    if pub["sector"] == INPUT_SECTOR_MAP.get(priv.get("sector", ""), ""):
        score += 40
    priv_rev = (priv.get("revenue_m") or 0) * 1e6
    pub_rev  = fin.get("revenue") or 0
    if priv_rev > 0 and pub_rev > 0:
        log_dist = abs(np.log10(pub_rev / priv_rev))
        score += max(0, 20 - log_dist * 8)
    priv_growth = (priv.get("rev_growth_pct") or 0) / 100
    pub_growth  = fin.get("rev_growth") or 0
    diff = abs(pub_growth - priv_growth)
    score += max(0, 20 - diff * 100)
    priv_gm = (priv.get("gross_margin_pct") or 0) / 100
    pub_gm  = fin.get("gross_margin") or 0
    diff = abs(pub_gm - priv_gm)
    score += max(0, 15 - diff * 60)
    if priv.get("geo") == pub.get("geo") or pub.get("geo") == "Global":
        score += 5
    return round(min(score, 100), 1)


# ── Models ────────────────────────────────────────────────────────────────────
class PrivateCompanyInput(BaseModel):
    company_name:      str
    sector:            str
    sub_sector:        Optional[str] = ""
    geo:               Optional[str] = "North America"
    stage:             Optional[str] = "Growth"
    revenue_m:         float
    ebitda_m:          Optional[float] = None
    gross_margin_pct:  Optional[float] = None
    net_income_m:      Optional[float] = None
    rev_growth_pct:    Optional[float] = None
    employees:         Optional[int]   = None
    max_comps:         Optional[int]   = 5

class ReportRequest(BaseModel):
    selectedCompany: Dict[str, Any]
    comps: List[Dict[str, Any]] = []
    privateCompany: Dict[str, Any] = {}
    multiples: Dict[str, Any] = {}
    implied: Dict[str, Any] = {}
    overall_range: Optional[Dict[str, Any]] = None
    sector_label: Optional[str] = ""
    comps_count: Optional[int] = None


class DeckRequest(ReportRequest):
    deckOptions: Dict[str, Any] = {}


class MARecommendationRequest(BaseModel):
    acquirer_ticker: Optional[str] = "CRM"
    target_ticker: Optional[str] = ""
    target_sector: Optional[str] = ""
    target_query: Optional[str] = ""
    max_results: Optional[int] = 5
    candidate_limit: Optional[int] = 120
    premium_pct: Optional[float] = 25.0
    cash_pct: Optional[float] = 35.0
    debt_pct: Optional[float] = 35.0
    stock_pct: Optional[float] = 30.0
    cost_synergy_pct: Optional[float] = 8.0
    revenue_synergy_pct: Optional[float] = 4.0


# ── Find Comps ────────────────────────────────────────────────────────────────
@app.post("/api/find-comps")
def find_comps(inp: PrivateCompanyInput):
    priv = inp.dict()
    target_sector = INPUT_SECTOR_MAP.get(inp.sector, "")
    candidates = [c for c in COMPANY_UNIVERSE if c["sector"] == target_sector]
    if not candidates:
        candidates = COMPANY_UNIVERSE

    scored = []
    for pub in candidates:
        fin = FINANCIALS_DB.get(pub["ticker"])
        if not fin:
            continue
        s = score_company(pub, priv, fin)
        scored.append({**pub, **fin, "match_score": s})

    scored.sort(key=lambda x: x["match_score"], reverse=True)
    top = scored[:inp.max_comps]

    if not top:
        return {"error": "No comparables found."}

    def stats(vals):
        vals = [v for v in vals if v is not None and v > 0]
        if not vals:
            return None
        return {
            "median": round(float(np.median(vals)), 2),
            "p25":    round(float(np.percentile(vals, 25)), 2),
            "p75":    round(float(np.percentile(vals, 75)), 2),
            "mean":   round(float(np.mean(vals)), 2),
            "min":    round(float(np.min(vals)), 2),
            "max":    round(float(np.max(vals)), 2),
        }

    multiples = {
        "ev_rev":    stats([c.get("ev_rev")    for c in top]),
        "ev_ebitda": stats([c.get("ev_ebitda") for c in top]),
        "ev_gp":     stats([c.get("ev_gp")     for c in top]),
        "pe":        stats([c.get("pe_ratio")   for c in top]),
    }

    rev_m    = inp.revenue_m
    ebitda_m = inp.ebitda_m
    gp_m     = rev_m * (inp.gross_margin_pct / 100) if inp.gross_margin_pct else None

    implied = {}
    if multiples["ev_rev"] and rev_m:
        m = multiples["ev_rev"]
        implied["ev_rev"] = {"method": "EV / Revenue", "low": round(m["p25"] * rev_m, 1), "mid": round(m["median"] * rev_m, 1), "high": round(m["p75"] * rev_m, 1), "multiple_used": m["median"], "label": f"${rev_m}M x {m['median']}x"}
    if multiples["ev_ebitda"] and ebitda_m and ebitda_m > 0:
        m = multiples["ev_ebitda"]
        implied["ev_ebitda"] = {"method": "EV / EBITDA", "low": round(m["p25"] * ebitda_m, 1), "mid": round(m["median"] * ebitda_m, 1), "high": round(m["p75"] * ebitda_m, 1), "multiple_used": m["median"], "label": f"${ebitda_m}M x {m['median']}x"}
    if multiples["ev_gp"] and gp_m and gp_m > 0:
        m = multiples["ev_gp"]
        implied["ev_gp"] = {"method": "EV / Gross Profit", "low": round(m["p25"] * gp_m, 1), "mid": round(m["median"] * gp_m, 1), "high": round(m["p75"] * gp_m, 1), "multiple_used": m["median"], "label": f"${round(gp_m,1)}M x {m['median']}x"}

    all_lows  = [v["low"]  for v in implied.values()]
    all_highs = [v["high"] for v in implied.values()]
    overall_range = {"low": round(min(all_lows), 1), "high": round(max(all_highs), 1)} if all_lows else None

    comps_out = []
    for c in top:
        comps_out.append({
            "ticker":       c["ticker"],
            "name":         c["name"],
            "sub":          c.get("sub", ""),
            "match_score":  c["match_score"],
            "market_cap_b": round(c["market_cap"] / 1e9, 2) if c.get("market_cap") else None,
            "ev_b":         round(c["ev"] / 1e9, 2)          if c.get("ev")         else None,
            "revenue_m":    round(c["revenue"] / 1e6, 1)     if c.get("revenue")    else None,
            "ebitda_m":     round(c["ebitda"] / 1e6, 1)      if c.get("ebitda")     else None,
            "ev_rev":       c.get("ev_rev"),
            "ev_ebitda":    c.get("ev_ebitda"),
            "ev_gp":        c.get("ev_gp"),
            "pe_ratio":     c.get("pe_ratio"),
            "rev_growth":   round(c["rev_growth"] * 100, 1)   if c.get("rev_growth")   else None,
            "gross_margin": round(c["gross_margin"] * 100, 1)  if c.get("gross_margin")  else None,
            "employees":    c.get("employees"),
        })

    return {
        "comps":         comps_out,
        "multiples":     multiples,
        "implied":       implied,
        "overall_range": overall_range,
        "sector_label":  SECTOR_LABEL_MAP.get(target_sector, inp.sector),
        "comps_count":   len(top),
    }


# ── M&A Recommendation Engine ────────────────────────────────────────────────
FALLBACK_MA_DEALS = [
    {
        "buyer": "Salesforce",
        "target": "Informatica",
        "value": "$8.0B",
        "sector": "Data Infrastructure",
        "status": "Announced",
        "date": "2025-05-27",
        "angle": "Strategic data cloud expansion across AI, governance, and integration.",
        "sourceName": "Salesforce Newsroom",
        "sourceUrl": "https://www.salesforce.com/news/press-releases/2025/05/27/salesforce-signs-definitive-agreement-to-acquire-informatica/",
    },
    {
        "buyer": "HPE",
        "target": "Juniper Networks",
        "value": "$14.0B",
        "sector": "AI Networking",
        "status": "Closed",
        "date": "2025-07-02",
        "angle": "Networking scale-up for AI-native enterprise and cloud infrastructure.",
        "sourceName": "HPE Newsroom",
        "sourceUrl": "https://www.hpe.com/us/en/newsroom/press-release/2025/07/hewlett-packard-enterprise-completes-acquisition-of-juniper-networks.html",
    },
    {
        "buyer": "Synopsys",
        "target": "Ansys",
        "value": "$35.0B",
        "sector": "Engineering Software",
        "status": "Closed",
        "date": "2025-07-17",
        "angle": "Silicon-to-systems platform consolidation across EDA and simulation.",
        "sourceName": "Synopsys Newsroom",
        "sourceUrl": "https://news.synopsys.com/2025-07-17-Synopsys-Completes-Acquisition-of-Ansys",
    },
    {
        "buyer": "Cisco",
        "target": "Splunk",
        "value": "$28.0B",
        "sector": "Observability",
        "status": "Closed",
        "date": "2024-03-18",
        "angle": "Security, observability, and data analytics consolidation.",
        "sourceName": "Cisco Newsroom",
        "sourceUrl": "https://newsroom.cisco.com/c/r/newsroom/en/us/a/y2024/m03/cisco-completes-acquisition-of-splunk.html",
    },
    {
        "buyer": "IBM",
        "target": "HashiCorp",
        "value": "$6.4B",
        "sector": "Cloud Infrastructure",
        "status": "Closed",
        "date": "2025-02-27",
        "angle": "Hybrid cloud automation and infrastructure lifecycle management.",
        "sourceName": "IBM Newsroom",
        "sourceUrl": "https://newsroom.ibm.com/2025-02-27-IBM-Completes-Acquisition-of-HashiCorp",
    },
    {
        "buyer": "Google",
        "target": "Wiz",
        "value": "$32.0B",
        "sector": "Cloud Security",
        "status": "Announced",
        "date": "2025-03-18",
        "angle": "Cloud-native security platform expansion for Google Cloud.",
        "sourceName": "Google Cloud Blog",
        "sourceUrl": "https://cloud.google.com/blog/products/identity-security/google-agreement-to-acquire-wiz",
    },
]


SECTOR_THEMES = {
    "public_unclassified": {
        "label": "All Public Companies",
        "themes": ["scale", "public-market liquidity", "strategic optionality"],
        "best_targets": ["saas_tech", "healthcare", "financial", "consumer", "industrials", "public_unclassified"],
    },
    "saas_tech": {
        "label": "Enterprise Software / SaaS",
        "themes": ["AI workflow automation", "data platform consolidation", "vertical SaaS expansion"],
        "best_targets": ["saas_tech", "financial", "healthcare"],
    },
    "healthcare": {
        "label": "Healthcare Technology",
        "themes": ["clinical workflow digitization", "connected diagnostics", "patient engagement"],
        "best_targets": ["healthcare", "saas_tech"],
    },
    "financial": {
        "label": "FinTech / Payments",
        "themes": ["embedded finance", "payment rails", "wealth and consumer finance distribution"],
        "best_targets": ["financial", "saas_tech", "consumer"],
    },
    "consumer": {
        "label": "Consumer / Retail",
        "themes": ["marketplace scale", "loyalty and commerce data", "category consolidation"],
        "best_targets": ["consumer", "financial", "saas_tech"],
    },
    "industrials": {
        "label": "Industrial Technology",
        "themes": ["automation", "smart infrastructure", "energy efficiency"],
        "best_targets": ["industrials", "saas_tech"],
    },
}

MA_SEED_COMPANIES = [
    {"ticker": "MSFT", "name": "Microsoft", "sector": "saas_tech", "sub": "Cloud / AI / Enterprise Software", "geo": "North America", "source": "Valence M&A seed"},
    {"ticker": "GOOGL", "name": "Alphabet", "sector": "saas_tech", "sub": "Search / Cloud / AI", "geo": "North America", "source": "Valence M&A seed"},
    {"ticker": "AAPL", "name": "Apple", "sector": "consumer", "sub": "Consumer Technology", "geo": "North America", "source": "Valence M&A seed"},
    {"ticker": "AMZN", "name": "Amazon", "sector": "consumer", "sub": "E-Commerce / Cloud", "geo": "North America", "source": "Valence M&A seed"},
    {"ticker": "ADBE", "name": "Adobe", "sector": "saas_tech", "sub": "Creative / Marketing Software", "geo": "North America", "source": "Valence M&A seed"},
    {"ticker": "CSCO", "name": "Cisco", "sector": "industrials", "sub": "Networking / Security", "geo": "North America", "source": "Valence M&A seed"},
    {"ticker": "PANW", "name": "Palo Alto Networks", "sector": "saas_tech", "sub": "Cybersecurity", "geo": "North America", "source": "Valence M&A seed"},
    {"ticker": "NET", "name": "Cloudflare", "sector": "saas_tech", "sub": "Network / Security", "geo": "North America", "source": "Valence M&A seed"},
    {"ticker": "MCD", "name": "McDonald's", "sector": "consumer", "sub": "Restaurants", "geo": "Global", "source": "Valence M&A seed"},
    {"ticker": "GIS", "name": "General Mills", "sector": "consumer", "sub": "Packaged Food", "geo": "North America", "source": "Valence M&A seed"},
    {"ticker": "KHC", "name": "Kraft Heinz", "sector": "consumer", "sub": "Packaged Food", "geo": "North America", "source": "Valence M&A seed"},
    {"ticker": "MDLZ", "name": "Mondelez International", "sector": "consumer", "sub": "Packaged Food", "geo": "Global", "source": "Valence M&A seed"},
]

MA_SEED_FINANCIALS = {
    "MSFT": {"market_cap": 3.7e12, "ev": 3.6e12, "revenue": 268e9, "ebitda": 140e9, "gross_profit": 185e9, "net_income": 95e9, "rev_growth": 0.15, "gross_margin": 0.69, "ev_rev": 13.4, "ev_ebitda": 25.7, "ev_gp": 19.5, "pe_ratio": 39},
    "GOOGL": {"market_cap": 2.2e12, "ev": 2.1e12, "revenue": 350e9, "ebitda": 125e9, "gross_profit": 205e9, "net_income": 92e9, "rev_growth": 0.13, "gross_margin": 0.59, "ev_rev": 6.0, "ev_ebitda": 16.8, "ev_gp": 10.2, "pe_ratio": 24},
    "AAPL": {"market_cap": 3.1e12, "ev": 3.2e12, "revenue": 390e9, "ebitda": 135e9, "gross_profit": 180e9, "net_income": 100e9, "rev_growth": 0.04, "gross_margin": 0.46, "ev_rev": 8.2, "ev_ebitda": 23.7, "ev_gp": 17.8, "pe_ratio": 31},
    "AMZN": {"market_cap": 2.3e12, "ev": 2.4e12, "revenue": 640e9, "ebitda": 115e9, "gross_profit": 310e9, "net_income": 60e9, "rev_growth": 0.11, "gross_margin": 0.48, "ev_rev": 3.8, "ev_ebitda": 20.9, "ev_gp": 7.7, "pe_ratio": 38},
    "ADBE": {"market_cap": 215e9, "ev": 210e9, "revenue": 21e9, "ebitda": 9e9, "gross_profit": 18e9, "net_income": 5.6e9, "rev_growth": 0.10, "gross_margin": 0.86, "ev_rev": 10.1, "ev_ebitda": 23.3, "ev_gp": 11.7, "pe_ratio": 38},
    "CSCO": {"market_cap": 210e9, "ev": 205e9, "revenue": 54e9, "ebitda": 18e9, "gross_profit": 34e9, "net_income": 11e9, "rev_growth": 0.02, "gross_margin": 0.63, "ev_rev": 3.8, "ev_ebitda": 11.4, "ev_gp": 6.0, "pe_ratio": 19},
    "PANW": {"market_cap": 125e9, "ev": 120e9, "revenue": 8.7e9, "ebitda": 2.2e9, "gross_profit": 6.5e9, "net_income": 1.1e9, "rev_growth": 0.16, "gross_margin": 0.75, "ev_rev": 13.8, "ev_ebitda": 54.5, "ev_gp": 18.5, "pe_ratio": 110},
    "NET": {"market_cap": 32e9, "ev": 31e9, "revenue": 1.7e9, "ebitda": 0.2e9, "gross_profit": 1.3e9, "net_income": 0.05e9, "rev_growth": 0.27, "gross_margin": 0.76, "ev_rev": 18.2, "ev_ebitda": 155.0, "ev_gp": 23.8, "pe_ratio": None},
    "MCD": {"market_cap": 220e9, "ev": 270e9, "revenue": 26e9, "ebitda": 14e9, "gross_profit": 15e9, "net_income": 8.5e9, "rev_growth": 0.04, "gross_margin": 0.58, "ev_rev": 10.4, "ev_ebitda": 19.3, "ev_gp": 18.0, "pe_ratio": 26},
    "GIS": {"market_cap": 35e9, "ev": 48e9, "revenue": 20e9, "ebitda": 4.5e9, "gross_profit": 7e9, "net_income": 2.5e9, "rev_growth": 0.01, "gross_margin": 0.35, "ev_rev": 2.4, "ev_ebitda": 10.7, "ev_gp": 6.9, "pe_ratio": 14},
    "KHC": {"market_cap": 40e9, "ev": 60e9, "revenue": 27e9, "ebitda": 6.2e9, "gross_profit": 9e9, "net_income": 3e9, "rev_growth": 0.01, "gross_margin": 0.33, "ev_rev": 2.2, "ev_ebitda": 9.7, "ev_gp": 6.7, "pe_ratio": 13},
    "MDLZ": {"market_cap": 90e9, "ev": 110e9, "revenue": 36e9, "ebitda": 7.5e9, "gross_profit": 14e9, "net_income": 4.8e9, "rev_growth": 0.04, "gross_margin": 0.39, "ev_rev": 3.1, "ev_ebitda": 14.7, "ev_gp": 7.9, "pe_ratio": 19},
}


YF_SECTOR_MAP = [
    ("saas_tech", ["software", "technology", "information technology", "communication services", "internet", "data", "cloud", "cyber", "semiconductor"]),
    ("healthcare", ["healthcare", "health care", "biotechnology", "medical", "pharmaceutical", "life sciences"]),
    ("financial", ["financial", "fintech", "bank", "insurance", "capital markets", "credit", "payments"]),
    ("consumer", ["consumer", "retail", "ecommerce", "restaurant", "apparel", "travel", "entertainment"]),
    ("industrials", ["industrial", "manufacturing", "machinery", "construction", "aerospace", "defense", "energy", "utilities", "materials"]),
]


def infer_public_sector(*values):
    text = " ".join([str(value or "") for value in values]).lower()
    for sector_key, needles in YF_SECTOR_MAP:
        if any(needle in text for needle in needles):
            return sector_key
    return "public_unclassified"


@lru_cache(maxsize=1)
def load_sec_public_company_universe():
    try:
        response = requests.get(TICKER_CIK_URL, headers=SEC_HEADERS, timeout=20)
        response.raise_for_status()
        data = response.json()
    except Exception:
        data = {}

    records = []
    seen = set()

    for item in data.values():
        ticker = str(item.get("ticker", "")).upper().strip()
        name = str(item.get("title", "")).strip()
        cik = str(item.get("cik_str", "")).zfill(10)
        if not ticker or ticker in seen:
            continue
        seen.add(ticker)
        records.append({
            "ticker": ticker,
            "name": name or ticker,
            "sector": "public_unclassified",
            "sub": "SEC public issuer",
            "geo": "Public markets",
            "cik": cik,
            "source": "SEC company_tickers.json",
        })

    for company in COMPANY_UNIVERSE:
        ticker = company["ticker"].upper()
        if ticker in seen:
            continue
        seen.add(ticker)
        records.append({**company, "source": "Valence curated universe"})

    for company in MA_SEED_COMPANIES:
        ticker = company["ticker"].upper()
        if ticker in seen:
            continue
        seen.add(ticker)
        records.append(company)

    return sorted(records, key=lambda row: row["ticker"])


@lru_cache(maxsize=512)
def fetch_market_financials(ticker: str):
    ticker = (ticker or "").upper().strip()
    if not ticker:
        return None

    curated = FINANCIALS_DB.get(ticker)
    if curated:
        return curated

    seeded = MA_SEED_FINANCIALS.get(ticker)
    if seeded:
        return seeded

    try:
        stock = yf.Ticker(ticker)
        info = stock.info or {}
        revenue = info.get("totalRevenue")
        ebitda = info.get("ebitda")
        gross_profit = info.get("grossProfits")
        ev = info.get("enterpriseValue")
        market_cap = info.get("marketCap")
        ev_rev = info.get("enterpriseToRevenue")
        ev_ebitda = info.get("enterpriseToEbitda")
        pe_ratio = info.get("trailingPE") or info.get("forwardPE")
        rev_growth = info.get("revenueGrowth")
        gross_margin = info.get("grossMargins")

        if not any([revenue, ev, market_cap, ev_rev, ev_ebitda]):
            return None

        if ev_rev is None and ev and revenue:
            ev_rev = safe_ratio(ev, revenue, None)
        if ev_ebitda is None and ev and ebitda and ebitda > 0:
            ev_ebitda = safe_ratio(ev, ebitda, None)
        ev_gp = safe_ratio(ev, gross_profit, None) if ev and gross_profit else None

        return {
            "market_cap": market_cap,
            "ev": ev,
            "revenue": revenue,
            "ebitda": ebitda,
            "gross_profit": gross_profit,
            "net_income": info.get("netIncomeToCommon") or info.get("netIncome"),
            "rev_growth": rev_growth,
            "gross_margin": gross_margin,
            "employees": info.get("fullTimeEmployees"),
            "pe_ratio": pe_ratio,
            "ev_rev": ev_rev,
            "ev_ebitda": ev_ebitda,
            "ev_gp": ev_gp,
            "yf_sector": info.get("sector"),
            "yf_industry": info.get("industry"),
            "long_business_summary": info.get("longBusinessSummary"),
        }
    except Exception:
        return None


def enrich_company_for_ma(company):
    ticker = company.get("ticker", "").upper()
    fin = fetch_market_financials(ticker)
    if not fin or not fin.get("ev"):
        return None

    sector = company.get("sector") or "public_unclassified"
    if sector == "public_unclassified":
        sector = infer_public_sector(fin.get("yf_sector"), fin.get("yf_industry"), company.get("name"))

    return {
        **company,
        **fin,
        "sector": sector,
        "sub": company.get("sub") if company.get("sub") != "SEC public issuer" else (fin.get("yf_industry") or "Public company"),
        "source": company.get("source") or ("Valence curated financials" if ticker in FINANCIALS_DB else "SEC universe + yfinance market data"),
    }


def company_by_ticker(ticker: str):
    ticker = (ticker or "").upper().strip()
    for company in COMPANY_UNIVERSE:
        if company.get("ticker") == ticker:
            fin = FINANCIALS_DB.get(ticker)
            if fin:
                return {**company, **fin}
    for company in load_sec_public_company_universe():
        if company.get("ticker") == ticker:
            return enrich_company_for_ma(company)
    return None


def safe_ratio(numerator, denominator, default=0):
    try:
        if denominator in [None, 0]:
            return default
        return float(numerator or 0) / float(denominator)
    except Exception:
        return default


def money_b(value):
    if value is None:
        return None
    return round(float(value) / 1e9, 2)


def pct(value):
    if value is None:
        return None
    return round(float(value) * 100, 1)


def deal_sector_hits(target):
    text = " ".join([
        target.get("sector", ""),
        target.get("sub", ""),
        SECTOR_THEMES.get(target.get("sector"), {}).get("label", ""),
    ]).lower()
    hits = []
    for deal in FALLBACK_MA_DEALS:
        haystack = " ".join([deal["sector"], deal["angle"], deal["buyer"], deal["target"]]).lower()
        if any(token and token in haystack for token in text.replace("/", " ").split()):
            hits.append(deal)
    return hits[:3]


SECTOR_ADJACENCY = {
    "saas_tech": {"saas_tech": 1.0, "financial": 0.68, "healthcare": 0.58, "industrials": 0.46, "consumer": 0.30, "public_unclassified": 0.36},
    "financial": {"financial": 1.0, "saas_tech": 0.66, "consumer": 0.52, "healthcare": 0.30, "industrials": 0.24, "public_unclassified": 0.32},
    "healthcare": {"healthcare": 1.0, "saas_tech": 0.58, "industrials": 0.34, "financial": 0.28, "consumer": 0.18, "public_unclassified": 0.30},
    "consumer": {"consumer": 1.0, "financial": 0.48, "saas_tech": 0.30, "industrials": 0.26, "healthcare": 0.18, "public_unclassified": 0.28},
    "industrials": {"industrials": 1.0, "saas_tech": 0.46, "healthcare": 0.34, "consumer": 0.26, "financial": 0.24, "public_unclassified": 0.30},
    "public_unclassified": {"public_unclassified": 0.42, "saas_tech": 0.36, "financial": 0.32, "healthcare": 0.30, "industrials": 0.30, "consumer": 0.28},
}


def sector_adjacency(acquirer, target):
    acq_sector = acquirer.get("sector") or "public_unclassified"
    target_sector = target.get("sector") or "public_unclassified"
    return SECTOR_ADJACENCY.get(acq_sector, {}).get(target_sector, 0.15)


def strategic_fit_score(acquirer, target):
    acq_sector = acquirer.get("sector")
    target_sector = target.get("sector")
    acq_theme = SECTOR_THEMES.get(acq_sector, {})
    target_text = f"{target.get('sub', '')} {SECTOR_THEMES.get(target_sector, {}).get('label', '')}".lower()
    adjacency = sector_adjacency(acquirer, target)

    score = 8 + adjacency * 32
    if target_sector in acq_theme.get("best_targets", []):
        score += 20
    if any(theme_word in target_text for theme in acq_theme.get("themes", []) for theme_word in theme.split()):
        score += 8
    return min(score, 55)


def build_merger_case(acquirer, target, req: MARecommendationRequest):
    acq_ev = float(acquirer.get("ev") or 0)
    target_ev = float(target.get("ev") or 0)
    acq_market_cap = float(acquirer.get("market_cap") or acq_ev or 0)
    target_revenue = float(target.get("revenue") or 0)
    target_ebitda = float(target.get("ebitda") or 0)
    acq_net_income = float(acquirer.get("net_income") or 0)
    target_net_income = float(target.get("net_income") or 0)

    premium = max(float(req.premium_pct or 0), 0) / 100
    purchase_ev = target_ev * (1 + premium)
    cash_pct = max(float(req.cash_pct or 0), 0) / 100
    debt_pct = max(float(req.debt_pct or 0), 0) / 100
    stock_pct = max(float(req.stock_pct or 0), 0) / 100
    total_mix = cash_pct + debt_pct + stock_pct
    if total_mix <= 0:
        cash_pct, debt_pct, stock_pct, total_mix = 0.35, 0.35, 0.30, 1
    cash_pct, debt_pct, stock_pct = cash_pct / total_mix, debt_pct / total_mix, stock_pct / total_mix

    cost_synergies = target_revenue * (max(float(req.cost_synergy_pct or 0), 0) / 100)
    revenue_synergies = target_revenue * (max(float(req.revenue_synergy_pct or 0), 0) / 100)
    incremental_ebitda = target_ebitda + cost_synergies + revenue_synergies * 0.25

    interest_rate = 0.065
    tax_rate = 0.24
    debt_cost = purchase_ev * debt_pct * interest_rate
    after_tax_income = (target_net_income + cost_synergies * (1 - tax_rate) + revenue_synergies * 0.12) - debt_cost * (1 - tax_rate)
    share_dilution = safe_ratio(purchase_ev * stock_pct, acq_market_cap)
    proforma_income = acq_net_income + after_tax_income
    eps_change = safe_ratio(proforma_income, 1 + share_dilution, acq_net_income) - acq_net_income
    eps_change_pct = safe_ratio(eps_change, abs(acq_net_income), 0) * 100 if acq_net_income else 0

    affordability = safe_ratio(purchase_ev, acq_ev, 9)
    affordability_score = max(0, 18 - affordability * 32)
    multiple_arbitrage = float(acquirer.get("ev_rev") or 0) - float(target.get("ev_rev") or 0)
    arbitrage_score = max(0, min(14, multiple_arbitrage * 2.2))
    growth_score = max(0, min(8, (float(target.get("rev_growth") or 0) - float(acquirer.get("rev_growth") or 0)) * 35 + 4))
    precedent_score = 8 if deal_sector_hits(target) else 3
    fit_score = strategic_fit_score(acquirer, target)
    model_score = max(0, min(12, eps_change_pct + 6))
    adjacency = sector_adjacency(acquirer, target)
    unrelated_penalty = 0
    if adjacency < 0.25:
        unrelated_penalty = 38
    elif adjacency < 0.40:
        unrelated_penalty = 22
    elif adjacency < 0.55:
        unrelated_penalty = 10

    total_score = round(max(5, min(99, fit_score + affordability_score + arbitrage_score + growth_score + precedent_score + model_score - unrelated_penalty)), 1)

    thesis = [
        f"{target['name']} adds {target.get('sub', 'category')} exposure to {acquirer['name']}.",
        f"Purchase EV is {affordability * 100:.1f}% of acquirer EV at a {req.premium_pct:.0f}% premium.",
        f"Modeled year-one EPS impact is {eps_change_pct:.1f}% using visible synergy assumptions.",
    ]
    if adjacency < 0.40:
        thesis.insert(0, "Strategic fit is weak: buyer and target operate in unrelated or low-adjacency sectors.")
    elif adjacency < 0.60:
        thesis.insert(0, "Strategic fit is moderate and depends on a clearly articulated expansion thesis.")
    else:
        thesis.insert(0, "Strategic fit is credible based on sector and business-model adjacency.")
    if multiple_arbitrage > 0:
        thesis.append(f"Acquirer trades {multiple_arbitrage:.1f}x EV/Revenue above target, supporting stock/deal currency logic.")

    return {
        "score": total_score,
        "acquirer": {
            "ticker": acquirer["ticker"],
            "name": acquirer["name"],
            "sector": SECTOR_THEMES.get(acquirer.get("sector"), {}).get("label", acquirer.get("sector")),
            "ev_b": money_b(acq_ev),
            "ev_rev": acquirer.get("ev_rev"),
            "rev_growth": pct(acquirer.get("rev_growth")),
        },
        "target": {
            "ticker": target["ticker"],
            "name": target["name"],
            "sector": SECTOR_THEMES.get(target.get("sector"), {}).get("label", target.get("sector")),
            "sub": target.get("sub"),
            "ev_b": money_b(target_ev),
            "revenue_b": money_b(target.get("revenue")),
            "ev_rev": target.get("ev_rev"),
            "rev_growth": pct(target.get("rev_growth")),
        },
        "model": {
            "purchase_ev_b": money_b(purchase_ev),
            "premium_pct": round(premium * 100, 1),
            "cash_pct": round(cash_pct * 100, 1),
            "debt_pct": round(debt_pct * 100, 1),
            "stock_pct": round(stock_pct * 100, 1),
            "cost_synergies_m": round(cost_synergies / 1e6, 1),
            "revenue_synergies_m": round(revenue_synergies / 1e6, 1),
            "incremental_ebitda_m": round(incremental_ebitda / 1e6, 1),
            "eps_change_pct": round(eps_change_pct, 1),
            "deal_size_to_acquirer_ev_pct": round(affordability * 100, 1),
        },
        "rationale": thesis,
        "risks": [
            "Needs diligence on customer overlap, churn, and product integration cost.",
            "Public-market multiples can shift before signing and closing.",
            "Regulatory review risk rises when buyer and target are close substitutes.",
        ],
        "precedents": deal_sector_hits(target),
    }


def company_matches_prefix(company: Dict[str, Any], search_text: str) -> bool:
    if not search_text:
        return True

    fields = [
        str(company.get("ticker", "")),
        str(company.get("name", "")),
        str(company.get("sub", "")),
    ]
    for field in fields:
        normalized = field.lower().strip()
        if normalized.startswith(search_text):
            return True
        if any(part.startswith(search_text) for part in normalized.replace("/", " ").replace("-", " ").split()):
            return True
    return False


@app.get("/api/ma/universe")
def ma_universe(limit: int = MA_UNIVERSE_LIMIT, search: str = ""):
    public_universe = load_sec_public_company_universe()
    search_text = (search or "").lower().strip()
    limit = max(1, min(int(limit or MA_UNIVERSE_LIMIT), 20000))
    companies = []

    for company in public_universe:
        if not company_matches_prefix(company, search_text):
            continue

        fin = FINANCIALS_DB.get(company["ticker"], {})
        companies.append({
            "ticker": company["ticker"],
            "name": company["name"],
            "sector": SECTOR_THEMES.get(company.get("sector"), {}).get("label", company.get("sector")),
            "raw_sector": company.get("sector"),
            "sub": company.get("sub", ""),
            "ev_b": money_b(fin.get("ev")),
            "market_cap_b": money_b(fin.get("market_cap")),
            "ev_rev": fin.get("ev_rev"),
            "rev_growth": pct(fin.get("rev_growth")),
            "has_financials": bool(fin),
            "source": company.get("source", "SEC company_tickers.json"),
        })
        if len(companies) >= limit:
            break

    return {
        "updatedAt": datetime.utcnow().isoformat() + "Z",
        "total_available": len(public_universe),
        "returned": len(companies),
        "coverage_note": "Universe comes from SEC company_tickers.json plus Valence curated records. Detailed valuation metrics are enriched on demand.",
        "companies": companies,
        "sectors": [{"key": key, "label": value["label"]} for key, value in SECTOR_THEMES.items()],
    }


@app.post("/api/ma/recommendations")
def ma_recommendations(req: MARecommendationRequest):
    acquirer = company_by_ticker(req.acquirer_ticker or "CRM")
    if not acquirer:
        raise HTTPException(status_code=404, detail="Acquirer not found or could not be enriched with market data.")

    max_results = max(1, min(int(req.max_results or 5), 5))
    candidate_limit = max(20, min(int(req.candidate_limit or 120), 400))
    target_sector = (req.target_sector or "").strip()
    target_query = (req.target_query or "").lower().strip()
    target_ticker = (req.target_ticker or "").upper().strip()
    raw_universe = load_sec_public_company_universe()

    curated_tickers = {company["ticker"] for company in COMPANY_UNIVERSE}
    candidates_raw = []
    seen = set()

    def candidate_matches(company):
        if company["ticker"] == acquirer["ticker"]:
            return False
        if target_query and not target_ticker:
            haystack = " ".join([
                company.get("ticker", ""),
                company.get("name", ""),
                company.get("sub", ""),
            ]).lower()
            if target_query not in haystack:
                return False
        if target_sector and company.get("sector") not in [target_sector, "public_unclassified"]:
            return False
        return True

    for company in COMPANY_UNIVERSE:
        if candidate_matches(company):
            seen.add(company["ticker"])
            candidates_raw.append(company)

    for company in raw_universe:
        ticker = company["ticker"]
        if ticker in seen or ticker in curated_tickers:
            continue
        if not candidate_matches(company):
            continue
        seen.add(ticker)
        candidates_raw.append(company)
        if len(candidates_raw) >= candidate_limit:
            break

    candidates = []

    if target_ticker:
        selected_target = company_by_ticker(target_ticker)
        if selected_target and selected_target["ticker"] != acquirer["ticker"]:
            candidates.append(build_merger_case(acquirer, selected_target, req))

    for company in candidates_raw[:candidate_limit]:
        if target_ticker and company["ticker"] == target_ticker:
            continue
        enriched = enrich_company_for_ma(company)
        if not enriched:
            continue
        if target_sector and enriched.get("sector") != target_sector:
            continue
        case = build_merger_case(acquirer, enriched, req)
        candidates.append(case)

    if target_ticker:
        candidates.sort(key=lambda item: (item["target"]["ticker"] != target_ticker, -item["score"]))
    else:
        candidates.sort(key=lambda item: item["score"], reverse=True)
    top = candidates[:max_results]

    return {
        "updatedAt": datetime.utcnow().isoformat() + "Z",
        "source": "SEC public-company universe + Valence curated data + yfinance market-data enrichment + official deal precedents",
        "universe_total": len(raw_universe),
        "screened_count": len(candidates_raw[:candidate_limit]),
        "financially_modeled_count": len(candidates),
        "assumptions": {
            "premium_pct": req.premium_pct,
            "cash_pct": req.cash_pct,
            "debt_pct": req.debt_pct,
            "stock_pct": req.stock_pct,
            "cost_synergy_pct": req.cost_synergy_pct,
            "revenue_synergy_pct": req.revenue_synergy_pct,
            "interest_rate": 6.5,
            "tax_rate": 24.0,
        },
        "recommendations": top,
        "dataArchitecture": [
            "SEC EDGAR company submissions and XBRL facts for public-company filings.",
            "Market-data providers for prices, enterprise value, ownership, and transaction feeds.",
            "News/event feeds such as GDELT for monitoring buyer intent, sector catalysts, and deal rumors.",
            "Private-company providers such as Crunchbase, PitchBook, CB Insights, or Apollo for smaller targets.",
        ],
        "modelLimitations": [
            "Public company list is broad SEC ticker coverage; smaller private-company scale still requires licensed private-company data.",
            "Non-curated tickers are enriched on demand, so very broad screens intentionally cap live financial lookups for speed.",
            "Recommendation score is a screening model, not fairness opinion or investment advice.",
            "A banker-grade output still needs management forecasts, debt schedule, purchase accounting, and integration plan.",
        ],
    }


@app.get("/api/market-intelligence/deals")
def latest_deals(limit: int = 6):
    return {
        "source": "official-fallback",
        "updatedAt": datetime.utcnow().isoformat() + "Z",
        "deals": FALLBACK_MA_DEALS[: max(1, min(limit, len(FALLBACK_MA_DEALS)))],
    }


# ── Generate PDF Report ───────────────────────────────────────────────────────
@app.post("/api/generate-report")
def generate_report(req: ReportRequest):
    company = req.selectedCompany
    comps   = req.comps
    private = req.privateCompany

    def to_float(value):
        try:
            return None if value in [None, "", "N/A", "NA"] else float(value)
        except (TypeError, ValueError):
            return None

    def adapt_company(raw):
        revenue_m = to_float(raw.get("revenue_m") or raw.get("revenue"))
        ebitda_m = to_float(raw.get("ebitda_m") or raw.get("ebitda"))
        ev_b = to_float(raw.get("ev_b"))
        market_cap_b = to_float(raw.get("market_cap_b"))
        ebitda_margin = (
            round(ebitda_m / revenue_m * 100, 1)
            if ebitda_m not in [None, 0] and revenue_m not in [None, 0]
            else raw.get("ebitdaMargin") or raw.get("ebitda_margin")
        )

        return {
            "companyName": raw.get("companyName") or raw.get("name") or raw.get("company") or "Company",
            "name": raw.get("name") or raw.get("companyName") or raw.get("company") or "Company",
            "ticker": raw.get("ticker") or raw.get("symbol") or "N/A",
            "sector": raw.get("sector") or raw.get("sub") or raw.get("industry") or "",
            "subSector": raw.get("sub") or raw.get("subSector") or raw.get("industry") or "",
            "description": raw.get("description") or raw.get("businessDescription") or "",
            "revenue": revenue_m,
            "ebitda": ebitda_m,
            "enterpriseValue": ev_b * 1000 if ev_b is not None else raw.get("enterpriseValue"),
            "marketCap": market_cap_b * 1000 if market_cap_b is not None else raw.get("marketCap"),
            "evRevenue": raw.get("ev_rev") or raw.get("evRevenue"),
            "evEbitda": raw.get("ev_ebitda") or raw.get("evEbitda"),
            "ebitdaMargin": ebitda_margin,
            "grossMargin": raw.get("gross_margin") or raw.get("grossMargin"),
            "revenueGrowth": raw.get("rev_growth") or raw.get("revenueGrowth"),
            "matchScore": raw.get("match_score") or raw.get("matchScore"),
            "marketDataDate": raw.get("marketDataDate") or "Latest available market/database data",
            "financialPeriod": raw.get("financialPeriod") or "Latest available fiscal year / LTM where available",
        }

    selected_company = adapt_company(company or {})
    comp_set = [adapt_company(c) for c in (comps or [company])]

    selected_ticker = selected_company.get("ticker")
    if selected_ticker:
        existing_tickers = {c.get("ticker") for c in comp_set}
        if selected_ticker not in existing_tickers:
            comp_set = [selected_company] + comp_set

    private_company = {
        "companyName": private.get("companyName") or private.get("company_name") or private.get("name") or "Target Company",
        "description": private.get("description") or private.get("businessDescription") or "Not provided",
        "sector": private.get("sector") or "N/A",
        "subSector": private.get("subSector") or private.get("sub_sector") or private.get("stage") or "",
        "geography": private.get("geography") or private.get("geo") or "N/A",
        "revenueModel": private.get("revenueModel") or private.get("revenue_model") or "Not provided",
        "customerType": private.get("customerType") or private.get("customer_type") or "Not provided",
        "revenue": private.get("revenue") or private.get("revenue_m"),
        "ebitda": private.get("ebitda") or private.get("ebitda_m"),
        "grossMargin": private.get("grossMargin") or private.get("gross_margin_pct"),
        "revenueGrowth": private.get("revenueGrowth") or private.get("rev_growth_pct"),
        "employees": private.get("employees") or private.get("employeeCount"),
    }

    ticker = selected_company.get("ticker", "company")
    output_path = os.path.join(tempfile.gettempdir(), f"valence-report-{ticker}.pdf")

    generate_banker_report(
        selected_company=selected_company,
        comps=comp_set,
        private_company=private_company,
        output_path=output_path,
    )

    return FileResponse(
        output_path,
        media_type="application/pdf",
        filename=f"valence-report-{ticker}.pdf",
    )


# ── Generate Consulting Slide Deck ────────────────────────────────────────────
@app.post("/api/generate-deck")
def generate_deck(req: DeckRequest):
    try:
        from pptx import Presentation
        from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR, PP_ALIGN
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from PIL import Image, ImageDraw, ImageFont
    except ImportError as exc:
        raise HTTPException(
            status_code=500,
            detail="PowerPoint export is not installed on the backend. Install python-pptx and retry.",
        ) from exc

    company = req.selectedCompany or {}
    comps = req.comps or []
    private = req.privateCompany or {}
    options = req.deckOptions or {}

    colorways = {
        "midnight": {
            "name": "Midnight Cyan",
            "bg": "07111F",
            "panel": "0F1B2D",
            "ink": "F8FAFC",
            "muted": "94A3B8",
            "accent": "22D3EE",
            "accent2": "7C3AED",
        },
        "boardroom": {
            "name": "Boardroom Blue",
            "bg": "F6F8FB",
            "panel": "FFFFFF",
            "ink": "102033",
            "muted": "5D6B7A",
            "accent": "1D4ED8",
            "accent2": "0F766E",
        },
        "emerald": {
            "name": "Emerald Slate",
            "bg": "071A16",
            "panel": "0D2B24",
            "ink": "ECFDF5",
            "muted": "A7F3D0",
            "accent": "10B981",
            "accent2": "F59E0B",
        },
        "plum": {
            "name": "Plum Strategy",
            "bg": "170923",
            "panel": "251137",
            "ink": "FAF5FF",
            "muted": "D8B4FE",
            "accent": "A855F7",
            "accent2": "22D3EE",
        },
    }

    theme = colorways.get(options.get("colorway"), colorways["midnight"])
    component_style = options.get("componentStyle", "strategy")
    density = options.get("density", "balanced")
    instructions = str(options.get("instructions") or "").strip()
    selected_sections = options.get("sections") or {}

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def rgb(hex_value):
        clean = str(hex_value).replace("#", "")
        return RGBColor(int(clean[0:2], 16), int(clean[2:4], 16), int(clean[4:6], 16))

    def fmt_m(v):
        if v is None: return "N/A"
        try: return f"${float(v):,.0f}M"
        except Exception: return "N/A"

    def fmt_b(v):
        if v is None: return "N/A"
        try: return f"${float(v):,.2f}B"
        except Exception: return "N/A"

    def fmt_x(v):
        if v is None: return "N/A"
        try: return f"{float(v):.1f}x"
        except Exception: return "N/A"

    def fmt_pct(v):
        if v is None: return "N/A"
        try: return f"{float(v):.1f}%"
        except Exception: return "N/A"

    def safe_text(value, fallback="N/A"):
        if value is None or value == "":
            return fallback
        return str(value)

    def hex_to_tuple(hex_value):
        clean = str(hex_value).replace("#", "")
        return tuple(int(clean[i:i + 2], 16) for i in (0, 2, 4))

    def fit_font_size(text, width, height, base_size, minimum=9):
        text_len = len(str(text or ""))
        capacity = max(18, int(width * height * 42))
        if text_len <= capacity:
            return base_size
        scale = capacity / max(text_len, 1)
        return max(minimum, int(base_size * max(0.62, scale)))

    def make_visual_panel(kind, title, points=None, width=1440, height=900):
        bg = hex_to_tuple(theme["panel"])
        accent = hex_to_tuple(theme["accent"])
        accent2 = hex_to_tuple(theme["accent2"])
        ink = hex_to_tuple(theme["ink"])
        muted = hex_to_tuple(theme["muted"])

        img = Image.new("RGB", (width, height), bg)
        draw = ImageDraw.Draw(img, "RGBA")

        for y in range(height):
            blend = y / max(height - 1, 1)
            r = int(bg[0] * (1 - blend) + max(accent2[0] - 20, 0) * blend)
            g = int(bg[1] * (1 - blend) + max(accent2[1] - 20, 0) * blend)
            b = int(bg[2] * (1 - blend) + max(accent2[2] - 20, 0) * blend)
            draw.line([(0, y), (width, y)], fill=(r, g, b, 255))

        for i in range(-height, width, 110):
            draw.line([(i, height), (i + height, 0)], fill=(*accent, 34), width=3)
        for i in range(11):
            x = 90 + i * ((width - 180) / 10)
            draw.line([(x, 120), (x, height - 120)], fill=(*muted, 28), width=1)
        for i in range(6):
            y = 160 + i * ((height - 280) / 5)
            draw.line([(90, y), (width - 90, y)], fill=(*muted, 28), width=1)

        values = points or [0.25, 0.42, 0.34, 0.58, 0.52, 0.71, 0.64]
        coords = []
        for idx, value in enumerate(values):
            x = 115 + idx * ((width - 230) / max(len(values) - 1, 1))
            y = height - 155 - (height - 340) * max(0, min(float(value), 1))
            coords.append((x, y))

        if kind == "bars":
            bar_w = (width - 260) / max(len(values), 1) * 0.62
            for idx, value in enumerate(values):
                x = 130 + idx * ((width - 260) / max(len(values), 1))
                bar_h = (height - 330) * max(0.05, min(float(value), 1))
                bar_color = accent if idx % 2 == 0 else accent2
                draw.rounded_rectangle(
                    [x, height - 145 - bar_h, x + bar_w, height - 145],
                    radius=18,
                    fill=(*bar_color, 210),
                )
        else:
            for offset, alpha in [(18, 42), (9, 80), (0, 230)]:
                shifted = [(x, y + offset) for x, y in coords]
                draw.line(shifted, fill=(*accent, alpha), width=10 if offset else 6, joint="curve")
            for x, y in coords:
                draw.ellipse([x - 13, y - 13, x + 13, y + 13], fill=(*accent2, 245), outline=(*ink, 190), width=3)

        font_title = ImageFont.load_default()
        draw.rounded_rectangle([70, 56, width - 70, 122], radius=24, fill=(*bg, 170), outline=(*accent, 130), width=2)
        draw.text((96, 78), title.upper(), font=font_title, fill=(*ink, 240))
        draw.text((96, height - 88), "Valence generated market visual", font=font_title, fill=(*muted, 190))

        buffer = io.BytesIO()
        img.save(buffer, format="PNG")
        buffer.seek(0)
        return buffer

    def add_visual(slide, kind, title, x, y, w, h, points=None):
        panel = make_visual_panel(kind, title, points=points)
        pic = slide.shapes.add_picture(panel, Inches(x), Inches(y), Inches(w), Inches(h))
        return pic

    def add_box(slide, x, y, w, h, fill=None, line=None, radius=False):
        shape_type = 5 if radius else 1
        box = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
        box.fill.solid()
        box.fill.fore_color.rgb = rgb(fill or theme["panel"])
        box.line.color.rgb = rgb(line or theme["accent"])
        box.line.transparency = 55
        return box

    def add_text(slide, text, x, y, w, h, size=18, bold=False, color=None, align=None, fit=True):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.clear()
        tf.margin_left = Inches(0.04)
        tf.margin_right = Inches(0.04)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        if fit:
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        p = tf.paragraphs[0]
        p.text = str(text)
        if align:
            p.alignment = align
        run = p.runs[0]
        run.font.name = "Aptos"
        run.font.size = Pt(fit_font_size(text, w, h, size) if fit else size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color or theme["ink"])
        return tb

    def add_bullets(slide, bullets, x, y, w, h, size=16, color=None):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.04)
        tf.margin_bottom = Inches(0.04)
        total_text = " ".join(str(b) for b in bullets if b)
        bullet_size = fit_font_size(total_text, w, h, size, minimum=10)
        for index, bullet in enumerate([b for b in bullets if b]):
            p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
            p.text = f"• {str(bullet)}"
            p.level = 0
            p.font.name = "Aptos"
            p.font.size = Pt(bullet_size)
            p.font.color.rgb = rgb(color or theme["ink"])
            p.space_after = Pt(8)
        return tb

    def add_header(slide, title, subtitle=""):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = rgb(theme["bg"])
        add_text(slide, "VALENCE", 0.55, 0.25, 1.4, 0.28, 9, True, theme["accent"])
        add_text(slide, title, 0.55, 0.58, 8.8, 0.55, 27, True)
        if subtitle:
            add_text(slide, subtitle, 0.58, 1.12, 8.4, 0.28, 11, False, theme["muted"])
        add_box(slide, 11.55, 0.28, 1.15, 0.08, theme["accent"], theme["accent"])

    def add_footer(slide, page):
        add_text(slide, f"{theme['name']} | {component_style.title()} layout", 0.55, 7.12, 4.2, 0.22, 8, False, theme["muted"])
        add_text(slide, str(page).zfill(2), 12.25, 7.08, 0.45, 0.24, 9, True, theme["accent"], PP_ALIGN.RIGHT)

    def metric_card(slide, label, value, x, y, w=2.2, accent=None):
        add_box(slide, x, y, w, 0.92, theme["panel"], accent or theme["accent"], True)
        add_text(slide, label.upper(), x + 0.14, y + 0.12, w - 0.28, 0.18, 7, True, theme["muted"])
        add_text(slide, value, x + 0.14, y + 0.36, w - 0.28, 0.32, 18, True, accent or theme["ink"])

    ticker = safe_text(company.get("ticker") or company.get("symbol"), "N/A")
    name = safe_text(company.get("name") or company.get("companyName") or company.get("company"), ticker)
    private_name = safe_text(private.get("company_name") or private.get("companyName") or private.get("name"), "Target Company")
    sector = safe_text(req.sector_label or private.get("sector") or company.get("sector"), "Selected sector")
    valuation_range = req.overall_range or {}
    top_comps = comps[:8]
    median_ev_rev = (req.multiples or {}).get("ev_rev", {}).get("median")
    median_ev_ebitda = (req.multiples or {}).get("ev_ebitda", {}).get("median")

    def normalized_points(values, fallback=None):
        nums = []
        for value in values:
            try:
                if value is not None:
                    nums.append(float(value))
            except Exception:
                continue
        if not nums:
            return fallback or [0.18, 0.34, 0.29, 0.55, 0.47, 0.74, 0.66]
        low = min(nums)
        high = max(nums)
        spread = high - low or 1
        return [0.16 + ((value - low) / spread) * 0.68 for value in nums]

    slide_num = 1

    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(theme["bg"])
    add_box(slide, 0.0, 0.0, 13.333, 7.5, theme["bg"], theme["bg"])
    add_box(slide, 8.55, 0.0, 4.8, 7.5, theme["panel"], theme["panel"])
    add_visual(
        slide,
        "line",
        "Comparable valuation signal",
        0.72,
        4.05,
        6.95,
        2.42,
        normalized_points([comp.get("ev_rev") for comp in top_comps]),
    )
    add_text(slide, "Consulting Competition Deck", 0.72, 0.72, 4.6, 0.32, 13, True, theme["accent"])
    add_text(slide, f"{private_name} comparable company analysis", 0.72, 1.2, 7.0, 1.18, 36, True)
    add_text(slide, "Valuation, peer benchmarking, strategic implications, and judge-ready talking points for a consulting team.", 0.78, 2.58, 6.35, 0.72, 17, False, theme["muted"])
    metric_card(slide, "Lead public comp", f"{name} ({ticker})", 8.95, 1.05, 3.55, theme["accent"])
    metric_card(slide, "Peer set", f"{len(top_comps) or len(comps)} companies", 8.95, 2.18, 3.55, theme["accent2"])
    metric_card(slide, "Implied EV range", f"{fmt_m(valuation_range.get('low'))} - {fmt_m(valuation_range.get('high'))}", 8.95, 3.31, 3.55, theme["accent"])
    metric_card(slide, "Median EV / Revenue", fmt_x(median_ev_rev), 8.95, 4.44, 3.55, theme["accent2"])
    add_text(slide, "Generated by Valence", 0.78, 6.94, 2.6, 0.25, 10, True, theme["muted"])
    add_footer(slide, slide_num)
    slide_num += 1

    if selected_sections.get("overview", True):
        slide = prs.slides.add_slide(blank)
        add_header(slide, "Executive answer", f"What the team should say in the first 90 seconds.")
        thesis = [
            f"{private_name} screens against {sector} comps, with {name} as the closest benchmark.",
            f"The selected peer set implies an enterprise value range of {fmt_m(valuation_range.get('low'))} to {fmt_m(valuation_range.get('high'))}.",
            "The story should balance comparability with caveats around growth, margin profile, and scale.",
            "Competition framing: defend the peer set, triangulate valuation, then translate the math into strategic options.",
        ]
        add_box(slide, 0.72, 1.42, 6.18, 4.42, theme["panel"], theme["accent"], True)
        add_bullets(slide, thesis, 0.95, 1.72, 5.62, 3.62, 16)
        metric_card(slide, "Target revenue", fmt_m(private.get("revenue_m") or private.get("revenue")), 7.25, 1.55, 2.45, theme["accent"])
        metric_card(slide, "Target EBITDA", fmt_m(private.get("ebitda_m") or private.get("ebitda")), 10.0, 1.55, 2.45, theme["accent2"])
        metric_card(slide, "Revenue growth", fmt_pct(private.get("rev_growth_pct") or private.get("revenueGrowth")), 7.25, 2.8, 2.45, theme["accent2"])
        metric_card(slide, "Gross margin", fmt_pct(private.get("gross_margin_pct") or private.get("grossMargin")), 10.0, 2.8, 2.45, theme["accent"])
        add_visual(
            slide,
            "bars",
            "Target profile snapshot",
            7.25,
            4.05,
            5.22,
            1.72,
            normalized_points([
                private.get("revenue_m") or private.get("revenue"),
                private.get("ebitda_m") or private.get("ebitda"),
                private.get("rev_growth_pct") or private.get("revenueGrowth"),
                private.get("gross_margin_pct") or private.get("grossMargin"),
            ]),
        )
        add_box(slide, 7.25, 5.98, 5.2, 0.62, theme["panel"], theme["accent"], True)
        add_text(slide, "Use as a briefing pack: narrative, valuation evidence, peer defense, risks, and workplan.", 7.46, 6.09, 4.75, 0.32, 12, False, theme["ink"])
        add_footer(slide, slide_num)
        slide_num += 1

    if selected_sections.get("valuation", True):
        slide = prs.slides.add_slide(blank)
        add_header(slide, "Valuation bridge", "Multiple-based range with a consulting-style recommendation frame.")
        methods = list((req.implied or {}).values())
        if not methods:
            methods = [
                {"method": "EV / Revenue", "low": None, "mid": None, "high": None, "label": "Add target revenue and peer multiples"},
            ]
        add_visual(
            slide,
            "line",
            "Valuation bridge visual",
            0.75,
            4.9,
            5.15,
            1.35,
            normalized_points([method.get("mid") for method in methods]),
        )
        for idx, method in enumerate(methods[:4]):
            y = 1.52 + idx * 0.9
            add_box(slide, 0.75, y, 11.85, 0.68, theme["panel"], theme["accent" if idx % 2 == 0 else "accent2"], True)
            add_text(slide, safe_text(method.get("method"), "Valuation method"), 0.98, y + 0.14, 2.25, 0.26, 13, True)
            add_text(slide, safe_text(method.get("label"), "Multiple x metric"), 3.35, y + 0.16, 2.4, 0.24, 11, False, theme["muted"])
            add_text(slide, fmt_m(method.get("low")), 6.25, y + 0.14, 1.25, 0.26, 14, True, theme["muted"], PP_ALIGN.RIGHT)
            add_text(slide, fmt_m(method.get("mid")), 8.0, y + 0.09, 1.5, 0.32, 19, True, theme["accent"], PP_ALIGN.CENTER)
            add_text(slide, fmt_m(method.get("high")), 10.05, y + 0.14, 1.25, 0.26, 14, True, theme["ink"], PP_ALIGN.RIGHT)
        add_box(slide, 6.25, 5.22, 6.35, 0.78, theme["accent"], theme["accent"], True)
        add_text(slide, f"Overall indicated enterprise value: {fmt_m(valuation_range.get('low'))} - {fmt_m(valuation_range.get('high'))}", 6.52, 5.4, 5.8, 0.32, 17, True, theme["bg"], PP_ALIGN.CENTER)
        add_footer(slide, slide_num)
        slide_num += 1

    if selected_sections.get("comps", True):
        slide = prs.slides.add_slide(blank)
        add_header(slide, "Peer landscape", "Selected public companies ranked by fit and valuation relevance.")
        headers = ["Ticker", "Company", "Match", "EV/Rev", "EV/EBITDA", "Growth", "GM"]
        col_x = [0.72, 1.65, 5.55, 6.48, 7.62, 9.0, 10.0]
        widths = [0.72, 3.45, 0.78, 0.78, 1.0, 0.74, 0.74]
        for i, header in enumerate(headers):
            add_text(slide, header.upper(), col_x[i], 1.42, widths[i], 0.2, 8, True, theme["accent"])
        for idx, comp in enumerate(top_comps[:7]):
            y = 1.84 + idx * 0.62
            add_box(slide, 0.58, y - 0.08, 10.94, 0.48, theme["panel"], theme["panel"], True)
            values = [
                safe_text(comp.get("ticker") or comp.get("symbol")),
                safe_text(comp.get("name") or comp.get("company")),
                fmt_pct(comp.get("match_score")),
                fmt_x(comp.get("ev_rev")),
                fmt_x(comp.get("ev_ebitda")),
                fmt_pct(comp.get("rev_growth")),
                fmt_pct(comp.get("gross_margin")),
            ]
            for i, value in enumerate(values):
                add_text(slide, value, col_x[i], y, widths[i], 0.22, 10 if i else 11, i in [0, 2], theme["ink" if i != 0 else "accent"])
        add_box(slide, 0.72, 6.0, 11.55, 0.62, theme["accent2"], theme["accent2"], True)
        add_text(slide, "Peer defense: explain why each company belongs, then acknowledge scale, growth, and margin differences.", 1.0, 6.13, 10.95, 0.3, 13, True, theme["bg"], PP_ALIGN.CENTER)
        add_footer(slide, slide_num)
        slide_num += 1

    if selected_sections.get("multiples", True):
        slide = prs.slides.add_slide(blank)
        add_header(slide, "Trading multiples dashboard", "The evidence judges will expect to see before the recommendation.")
        multiple_rows = [
            ("EV / Revenue", (req.multiples or {}).get("ev_rev")),
            ("EV / EBITDA", (req.multiples or {}).get("ev_ebitda")),
            ("EV / Gross Profit", (req.multiples or {}).get("ev_gp")),
            ("P / E", (req.multiples or {}).get("pe")),
        ]
        add_visual(
            slide,
            "bars",
            "Multiple dispersion",
            0.75,
            5.72,
            4.25,
            0.9,
            normalized_points([(data or {}).get("median") for _, data in multiple_rows]),
        )
        for idx, (label, data) in enumerate(multiple_rows):
            x = 0.75 + (idx % 2) * 6.0
            y = 1.55 + (idx // 2) * 2.25
            add_box(slide, x, y, 5.55, 1.7, theme["panel"], theme["accent" if idx % 2 == 0 else "accent2"], True)
            add_text(slide, label, x + 0.25, y + 0.22, 2.7, 0.32, 17, True, theme["accent" if idx % 2 == 0 else "accent2"])
            add_text(slide, f"Median {fmt_x((data or {}).get('median'))}", x + 0.25, y + 0.68, 2.5, 0.34, 22, True)
            add_text(slide, f"IQR: {fmt_x((data or {}).get('p25'))} - {fmt_x((data or {}).get('p75'))}", x + 3.05, y + 0.72, 2.05, 0.28, 14, False, theme["muted"])
            add_text(slide, f"Range: {fmt_x((data or {}).get('min'))} - {fmt_x((data or {}).get('max'))}", x + 3.05, y + 1.05, 2.05, 0.24, 12, False, theme["muted"])
        add_box(slide, 5.25, 5.85, 7.05, 0.64, theme["panel"], theme["accent"], True)
        add_text(slide, "Talk track: median is the base case; IQR is the defensible valuation corridor.", 5.55, 5.99, 6.4, 0.28, 13, False, theme["ink"])
        add_footer(slide, slide_num)
        slide_num += 1

    slide = prs.slides.add_slide(blank)
    add_header(slide, "Competition playbook", "How to turn the analysis into a winning room narrative.")
    workplan = [
        "Open with the answer: valuation range, recommended midpoint, and why the peer set is credible.",
        "Show the target profile before the comps so judges understand what the benchmark must explain.",
        "Use the multiples dashboard to defend method selection and avoid over-relying on one metric.",
        "Translate valuation into strategic implications: pricing, funding, acquisition logic, or growth priorities.",
        "Close with risks and diligence questions instead of pretending the model is perfect.",
    ]
    if density == "detailed":
        workplan.append("Assign one teammate to own valuation, one to own market story, and one to own Q&A defense.")
    add_box(slide, 0.75, 1.5, 6.65, 4.85, theme["panel"], theme["accent"], True)
    add_bullets(slide, workplan, 0.98, 1.82, 6.12, 3.98, 15)
    add_visual(slide, "line", "Presentation arc", 7.72, 1.5, 4.65, 1.8, [0.2, 0.38, 0.52, 0.7, 0.86])
    add_box(slide, 7.72, 3.62, 4.65, 2.72, theme["panel"], theme["accent2"], True)
    add_text(slide, "Custom deck brief", 8.02, 3.9, 2.8, 0.32, 15, True, theme["accent2"])
    add_text(slide, instructions or "No custom instructions were provided. The deck uses Valence's default consulting competition structure.", 8.02, 4.34, 3.92, 1.52, 13, False, theme["ink"])
    add_footer(slide, slide_num)
    slide_num += 1

    if selected_sections.get("market", False):
        slide = prs.slides.add_slide(blank)
        add_header(slide, "Appendix: diligence questions", "Questions the team should be ready to answer.")
        add_bullets(slide, [
            "Which peers should be excluded because of business model, scale, geography, or profitability differences?",
            "How sensitive is valuation to revenue growth, EBITDA margin, and gross margin assumptions?",
            "What strategic buyer or investor would pay a premium, and what synergy logic supports it?",
            "What data would change the recommendation if the team had another week of diligence?",
        ], 0.85, 1.55, 10.8, 3.8, 19)
        add_footer(slide, slide_num)
        slide_num += 1

    filename_ticker = "".join(ch for ch in ticker.upper() if ch.isalnum() or ch in ["_", "-"]) or "COMPANY"
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="valence-deck-{filename_ticker}.pptx"'}
    )


# ── Health & Sectors ──────────────────────────────────────────────────────────
@app.get("/api/health")
def health():
    return {"status": "ok", "companies_in_db": len(FINANCIALS_DB)}

@app.get("/api/sectors")
def sectors():
    return {"sectors": list(INPUT_SECTOR_MAP.keys())}


# ── SEC EDGAR Helpers ─────────────────────────────────────────────────────────
@lru_cache(maxsize=1)
def load_ticker_to_cik_map():
    response = requests.get(TICKER_CIK_URL, headers=SEC_HEADERS, timeout=20)
    response.raise_for_status()
    data = response.json()
    ticker_map = {}
    for item in data.values():
        ticker = item["ticker"].upper()
        cik = str(item["cik_str"]).zfill(10)
        ticker_map[ticker] = cik
    return ticker_map

def sec_get_json(url):
    response = requests.get(url, headers=SEC_HEADERS, timeout=20)
    if response.status_code == 404:
        raise HTTPException(status_code=404, detail="SEC data not found.")
    response.raise_for_status()
    return response.json()

def build_filing_url(cik, accession_number, primary_document):
    cik_no_leading_zeros = str(int(cik))
    accession_no_dashes = accession_number.replace("-", "")
    return (
        f"https://www.sec.gov/Archives/edgar/data/"
        f"{cik_no_leading_zeros}/{accession_no_dashes}/{primary_document}"
    )

def get_recent_filings(submissions, cik):
    recent = submissions.get("filings", {}).get("recent", {})
    forms = recent.get("form", [])
    filing_dates = recent.get("filingDate", [])
    report_dates = recent.get("reportDate", [])
    accession_numbers = recent.get("accessionNumber", [])
    primary_docs = recent.get("primaryDocument", [])
    important_forms = {"10-K","10-Q","8-K","S-1","DEF 14A","10-K/A","10-Q/A","424B4"}
    filings = []
    for i, form in enumerate(forms):
        if form not in important_forms:
            continue
        accession = accession_numbers[i]
        primary_doc = primary_docs[i]
        filings.append({
            "form": form,
            "filingDate": filing_dates[i],
            "reportDate": report_dates[i] if i < len(report_dates) else None,
            "accessionNumber": accession,
            "document": primary_doc,
            "url": build_filing_url(cik, accession, primary_doc),
        })
        if len(filings) >= 25:
            break
    return filings

def latest_fact(companyfacts, possible_tags, unit="USD"):
    facts_root = companyfacts.get("facts", {}).get("us-gaap", {})
    for tag in possible_tags:
        concept = facts_root.get(tag)
        if not concept:
            continue
        unit_data = concept.get("units", {}).get(unit)
        if not unit_data:
            continue
        clean_items = [item for item in unit_data if item.get("val") is not None and item.get("filed")]
        if not clean_items:
            continue
        latest = sorted(clean_items, key=lambda x: (x.get("filed", ""), x.get("end", "")), reverse=True)[0]
        return {
            "tag":    tag,
            "label":  concept.get("label", tag),
            "value":  latest.get("val"),
            "filed":  latest.get("filed"),
            "period": latest.get("fp"),
            "form":   latest.get("form"),
            "end":    latest.get("end"),
        }
    return None

def get_financials(companyfacts):
    return {
        "revenue":           latest_fact(companyfacts, ["RevenueFromContractWithCustomerExcludingAssessedTax","SalesRevenueNet","Revenues"]),
        "netIncome":         latest_fact(companyfacts, ["NetIncomeLoss"]),
        "totalAssets":       latest_fact(companyfacts, ["Assets"]),
        "totalLiabilities":  latest_fact(companyfacts, ["Liabilities","LiabilitiesCurrent"]),
        "stockholdersEquity":latest_fact(companyfacts, ["StockholdersEquity","StockholdersEquityIncludingPortionAttributableToNoncontrollingInterest"]),
        "cash":              latest_fact(companyfacts, ["CashAndCashEquivalentsAtCarryingValue","CashCashEquivalentsRestrictedCashAndRestrictedCashEquivalents"]),
        "operatingIncome":   latest_fact(companyfacts, ["OperatingIncomeLoss"]),
    }

def safe_float(value):
    try:
        if value is None: return None
        return float(value)
    except: return None

def fast_info_get(fast_info, key):
    try: return fast_info[key]
    except:
        try: return fast_info.get(key)
        except: return None

def get_stock_data(ticker):
    try:
        stock = yf.Ticker(ticker)
        fast_info = stock.fast_info
        one_year_history = stock.history(period="1y")
        five_day_history = stock.history(period="5d")
        current_price = safe_float(fast_info_get(fast_info, "last_price"))
        if current_price is None and not five_day_history.empty:
            current_price = safe_float(five_day_history["Close"].dropna().iloc[-1])
        year_high = year_low = None
        if not one_year_history.empty:
            year_high = safe_float(one_year_history["High"].max())
            year_low  = safe_float(one_year_history["Low"].min())
        return {
            "ticker":        ticker,
            "currentPrice":  current_price,
            "previousClose": safe_float(fast_info_get(fast_info, "previous_close")),
            "dayHigh":       safe_float(fast_info_get(fast_info, "day_high")),
            "dayLow":        safe_float(fast_info_get(fast_info, "day_low")),
            "yearHigh":      year_high,
            "yearLow":       year_low,
            "marketCap":     safe_float(fast_info_get(fast_info, "market_cap")),
            "currency":      fast_info_get(fast_info, "currency") or "USD",
        }
    except Exception as e:
        return {"ticker": ticker, "error": str(e)}


# ── Company Detail ────────────────────────────────────────────────────────────
@app.get("/api/company/{ticker}")
def get_company_detail(ticker: str):
    ticker = ticker.upper().strip()
    ticker_map = load_ticker_to_cik_map()
    if ticker not in ticker_map:
        raise HTTPException(status_code=404, detail=f"No SEC CIK found for ticker {ticker}.")
    cik = ticker_map[ticker]
    submissions  = sec_get_json(f"https://data.sec.gov/submissions/CIK{cik}.json")
    companyfacts = sec_get_json(f"https://data.sec.gov/api/xbrl/companyfacts/CIK{cik}.json")
    return {
        "company": {
            "name":        submissions.get("name"),
            "ticker":      ticker,
            "cik":         cik,
            "sic":         submissions.get("sic"),
            "sicDescription": submissions.get("sicDescription"),
            "exchanges":   submissions.get("exchanges", []),
            "entityType":  submissions.get("entityType"),
        },
        "stock":    get_stock_data(ticker),
        "financials": get_financials(companyfacts),
        "filings":  get_recent_filings(submissions, cik),
    }
